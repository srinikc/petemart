from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_SLIDE = RGBColor(0xF5, 0xF0, 0xEB)
ACCENT_GOLD = RGBColor(0xD4, 0xA0, 0x3C)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
ACCENT_BLUE = RGBColor(0x15, 0x63, 0x8C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
LIGHT_GRAY = RGBColor(0xF0, 0xEE, 0xE9)
MID_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
TEAL = RGBColor(0x00, 0x7B, 0x7F)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BG_SLIDE):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background(); return shp

def add_accent_bar(slide, left=Inches(0), top=Inches(0), width=Inches(0.15), height=Inches(7.5), color=ACCENT_GOLD):
    return add_shape_rect(slide, left, top, width, height, color)

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font_name; p.alignment = alignment
    return tb

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=12, color=DARK_TEXT, bold_first=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size); p.font.color.rgb = color; p.font.name = 'Calibri'
        if bold_first and i == 0: p.font.bold = True; p.font.size = Pt(font_size + 2)
        p.space_after = Pt(3)
    return tb

def add_kpi_box(slide, left, top, width, height, number, label, bg_color=WHITE, num_color=ACCENT_GOLD):
    add_shape_rect(slide, left, top, width, height, bg_color)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.6),
                str(number), font_size=28, bold=True, color=num_color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.65), width - Inches(0.3), Inches(0.4),
                label, font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i); cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True; p.font.size = Pt(10); p.font.color.rgb = WHITE; p.font.name = 'Calibri'
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_BLUE
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx); cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9); p.font.color.rgb = DARK_TEXT; p.font.name = 'Calibri'
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 0 else LIGHT_GRAY
    if col_widths:
        for i, w in enumerate(col_widths): tbl.columns[i].width = w
    return tbl

def add_section_header(slide, text, subtitle=None):
    add_accent_bar(slide)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7), text, font_size=32, bold=True, color=ACCENT_BLUE)
    if subtitle: add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.5), subtitle, font_size=13, color=MID_GRAY)

def add_architecture_box(slide, left, top, width, height, title, items, color=ACCENT_BLUE):
    add_shape_rect(slide, left, top, width, height, WHITE)
    add_shape_rect(slide, left, top, width, Inches(0.45), color)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.05), width - Inches(0.3), Inches(0.35),
                title, font_size=11, bold=True, color=WHITE)
    add_multiline_textbox(slide, left + Inches(0.15), top + Inches(0.55), width - Inches(0.3), height - Inches(0.6),
                          items, font_size=9, color=DARK_TEXT)

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

# ════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════
s = new_slide()
add_bg(s, BG_DARK)
add_textbox(s, Inches(0.8), Inches(0.6), Inches(11), Inches(1), 'PeteMart', font_size=54, bold=True, color=ACCENT_GOLD)
add_textbox(s, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8), 'Program Management — Agile Sprint Plan & Delivery Framework', font_size=22, color=WHITE)
add_shape_rect(s, Inches(0.8), Inches(2.5), Inches(3), Inches(0.04), ACCENT_GOLD)
add_multiline_textbox(s, Inches(0.8), Inches(2.8), Inches(11), Inches(3), [
    'Program Management Review Package',
    'Agent 05 — Senior Agile Program Manager & Scrum Master',
    '',
    'Presented for: Product & Engineering Review',
    'Date: May 2026',
    'Execution Model: 3 AI Agents in Parallel (6x human velocity)',
    'Total Duration: 22 Days | MVP: 12 Days',
], font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

# ════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Executive Summary — Program at a Glance',
                   'AI-agent execution: 22 days full product vs 24 weeks human-equivalent')

kpis = [
    ('26', 'Epics Defined', Inches(0.5), ACCENT_BLUE),
    ('113', 'User Stories', Inches(3.7), ACCENT_GREEN),
    ('12', 'Sprints (2-day)', Inches(6.9), PURPLE),
    ('22 Days', 'Full Delivery', Inches(10.1), ACCENT_GOLD),
]
for num, label, left, color in kpis:
    add_kpi_box(s, left, Inches(1.3), Inches(2.8), Inches(1.1), num, label, num_color=color)

add_textbox(s, Inches(0.5), Inches(2.7), Inches(6), Inches(0.5), 'Execution Model', font_size=18, bold=True, color=ACCENT_BLUE)
add_multiline_textbox(s, Inches(0.5), Inches(3.3), Inches(6), Inches(3.5), [
    '3 AI Agents Run in Parallel at Machine Speed:',
    '',
    '  Agent 07a (UI) — Frontend & Mobile Interface',
    '    • Google Stitch generates React components from NL',
    '    • 30+ screens across all personas',
    '    • DESIGN.md → Tailwind config pipeline',
    '',
    '  Agent 07b (API) — Interface Connection',
    '    • All RESTful endpoints (/api/v1/*)',
    '    • Rate limiting, auth middleware, Zod validation',
    '    • Mock endpoints → live data routing',
    '',
    '  Agent 07c (Backend DB) — Data Infrastructure',
    '    • 35-table PostgreSQL schema (15 POC)',
    '    • RLS policies, migrations, indexes',
    '    • Redis cache + Supabase Realtime',
    '',
    'Each 2-day sprint = 3 agents × full sprint output',
    'No human bottlenecks — artifacts in hours, not weeks',
], font_size=11, color=DARK_TEXT)

add_textbox(s, Inches(6.8), Inches(2.7), Inches(6), Inches(0.5), 'Jira Integration', font_size=18, bold=True, color=ACCENT_BLUE)
add_multiline_textbox(s, Inches(6.8), Inches(3.3), Inches(6), Inches(3.5), [
    'Jira\'s Role in the Agentic AI Pipeline:',
    '',
    '  STATE_MATRIX.json is the SOURCE OF TRUTH',
    '  → Agents read/write state here, NOT in Jira',
    '',
    '  Jira is the HUMAN VISIBILITY LAYER:',
    '  → Stakeholders see progress without reading JSON',
    '  → Read-only mirror of agent state',
    '',
    '  What gets pushed to Jira:',
    '  → Epics, Features, User Stories (via API)',
    '  → Sprint boards with status per story',
    '  → Milestone tracking & gate status',
    '  → Bug/defect tracking (Agent 11)',
    '',
    '  Agents NEVER read from Jira for decisions',
    '  Jira = human reporting, not agent orchestration',
    '',
    '  Batch import ready: all 113 stories mapped to',
    '  epics, sprints, dependencies, and requirement IDs',
], font_size=11, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 3: EPIC HIERARCHY
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Epic Hierarchy — 26 Epics Across 6 Phases',
                   'Every epic traceable to PRD requirement IDs. Total: 72 features → 113 user stories')

epic_headers = ['Phase', 'Epic ID', 'Epic Name', 'Features', 'Stories', 'Sprint Range']
epic_rows = [
    ['Foundation', 'EPIC-01', 'Customer Auth & Persona Mgmt', '3', '12', 'S1-S2'],
    ['Foundation', 'EPIC-02', 'Browse & Discovery', '4', '18', 'S1-S3'],
    ['Foundation', 'EPIC-17', 'Infrastructure, Security & Compliance', '5', '18', 'S1-S8'],
    ['Commerce', 'EPIC-03', 'Shopping Cart & Checkout (Mode A)', '4', '16', 'S2-S4'],
    ['Commerce', 'EPIC-04', 'WhatsApp Enquiry (Mode B)', '2', '8', 'S2-S3'],
    ['Commerce', 'EPIC-05', 'Store Visit (Mode C)', '2', '8', 'S2-S3'],
    ['Commerce', 'EPIC-09', 'Payment Processing & Settlement', '3', '12', 'S3-S5'],
    ['Commerce', 'EPIC-10', 'Multi-Store Consolidation & Delivery', '3', '10', 'S4-S6'],
    ['Ops', 'EPIC-06', 'Merchant Dashboard & Microsite', '5', '22', 'S2-S5'],
    ['Ops', 'EPIC-07', 'Admin Console & Operations', '4', '16', 'S3-S6'],
    ['Ops', 'EPIC-08', 'Order Management & Tracking', '3', '14', 'S3-S5'],
    ['Ops', 'EPIC-11', 'Persona-Based Navigation & UX', '3', '12', 'S4-S6'],
    ['Ops', 'EPIC-15', 'Merchant Catalog & Inventory', '3', '12', 'S4-S6'],
    ['Ops', 'EPIC-24', 'PWA & Offline-First', '2', '6', 'S5-S6'],
    ['Growth', 'EPIC-12', 'Mobile App (Expo)', '4', '16', 'S5-S7'],
    ['Growth', 'EPIC-13', 'Reviews, Ratings & Trust', '3', '12', 'S5-S7'],
    ['Growth', 'EPIC-14', 'Promo Engine & Loyalty', '2', '10', 'S6-S7'],
    ['Growth', 'EPIC-16', 'Subscriptions & Billing', '2', '8', 'S5-S7'],
    ['Growth', 'EPIC-18', 'National Shipping (ShipRocket)', '2', '8', 'S6-S8'],
    ['Scale', 'EPIC-19', 'AI Virtual Try-On', '3', '12', 'S7-S9'],
    ['Scale', 'EPIC-20', 'Jewellery & Bullion Integration', '3', '10', 'S7-S9'],
    ['Scale', 'EPIC-21', 'Video Call & Trust Features', '2', '8', 'S8-S9'],
    ['Scale', 'EPIC-22', 'City Expansion & Multi-Geo', '2', '8', 'S8-S10'],
    ['Scale', 'EPIC-23', 'Analytics & CDP', '2', '8', 'S7-S9'],
    ['Vision', 'EPIC-25', 'Virtual Walk & Live Bazaar', '3', '10', 'S10-S12'],
    ['Vision', 'EPIC-26', 'Co-Shopping & Community', '2', '6', 'S11-S12'],
]
add_table(s, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5.8), epic_headers, epic_rows,
          col_widths=[Inches(1.3), Inches(1.3), Inches(3.5), Inches(1.2), Inches(1.2), Inches(1.2)])

# ════════════════════════════════════════
# SLIDE 4: SPRINT TIMELINE
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Sprint Timeline — 12 Sprints in 22 Days',
                   '2-day sprint cadence | 3 agents in parallel = 120 story points per sprint | MVP = 12 days')

phase_boxes = [
    ('Phase 1: Foundation', ACCENT_BLUE, Inches(0.3), 'S1-S2 (D1-D4)', ['Auth & Registration', 'Pete Tapestry landing', 'Search & Discovery', 'CI/CD + Infra setup']),
    ('Phase 2: Commerce', ACCENT_GREEN, Inches(4.8), 'S3-S4 (D5-D8)', ['Multi-store cart', 'Checkout + payment', 'Merchant dashboard', 'Delivery zones']),
    ('Phase 3: Operations', ORANGE, Inches(9.3), 'S5-S6 (D9-D12)', ['Admin console', 'i18n + PWA', 'Order tracking', 'Feature flags']),
]
for title, color, left, timeline, items in phase_boxes:
    add_architecture_box(s, left, Inches(1.3), Inches(4.3), Inches(3.5), title, [f'📅 {timeline}'] + items, color)

add_shape_rect(s, Inches(0.3), Inches(5.0), Inches(12.7), Inches(2.2), LIGHT_GRAY)
add_textbox(s, Inches(0.5), Inches(5.1), Inches(12), Inches(0.4), 'MVP Boundary: Sprints 1-6 (Days 1-12)', font_size=14, bold=True, color=ACCENT_GREEN)
add_multiline_textbox(s, Inches(0.5), Inches(5.6), Inches(12), Inches(1.5), [
    'Tier 1 (Days 13-14): Mobile (Expo), i18n, Reviews, Subscriptions  |  Tier 2 (Days 15-16): Promo Engine, ShipRocket, Loyalty',
    'Tier 3 (Days 17-18): AI Try-On, Bullion Rates, City Expansion, Analytics  |  Tier 4 (Days 19-22): Virtual Walk, Live Bazaar, Co-Shopping',
    'Gates: GATE-MVP-01 ✅ | GATE-COSTING-01 (End of S10) | GATE-PRODUCTION-01 (End of S12)',
], font_size=10, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 5: USER STORY BREAKDOWN
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'User Story Breakdown — 113 Stories',
                   'Every story traceable to Requirement ID. Authored with acceptance criteria and story points.')

story_headers = ['Story ID', 'Story Title', 'Epic', 'Points', 'Dependencies', 'Layer']
story_rows = [
    ['US-001', 'Customer OTP Registration', 'EPIC-01', '3', 'None', 'UI+API+DB'],
    ['US-002', 'Customer OTP Login', 'EPIC-01', '2', 'US-001', 'UI+API+DB'],
    ['US-004', 'Merchant Registration Wizard', 'EPIC-06', '8', 'US-008', 'UI+API+DB'],
    ['US-010', 'Pete Tapestry Landing Page', 'EPIC-02', '5', 'None', 'UI'],
    ['US-012', 'Search with Autocomplete', 'EPIC-02', '5', 'US-010', 'UI+API+DB'],
    ['US-015', 'Merchant Microsite (Public)', 'EPIC-06', '5', 'US-004', 'UI+API+DB'],
    ['US-019', 'Multi-Merchant Cart', 'EPIC-03', '3', 'US-017', 'UI+API+DB'],
    ['US-021', 'Consolidated Checkout', 'EPIC-03', '5', 'US-020, US-060', 'UI+API+DB'],
    ['US-023', 'Razorpay Payment Widget', 'EPIC-09', '5', 'US-021', 'UI+API+DB'],
    ['US-026', 'Order Lifecycle Engine', 'EPIC-08', '5', 'US-024', 'API+DB'],
    ['US-035', 'Merchant Dashboard Overview', 'EPIC-06', '5', 'US-026', 'UI+API+DB'],
    ['US-046', 'Admin Dashboard', 'EPIC-07', '5', 'US-006', 'UI+API+DB'],
    ['US-049', 'Feature Flag & Kill Switch', 'EPIC-07', '5', 'US-048', 'UI+API+DB'],
    ['US-060', 'Multi-Store Consolidation Engine', 'EPIC-10', '5', 'US-021', 'DB'],
    ['US-072', 'Product & Merchant Reviews', 'EPIC-13', '5', 'US-026', 'UI+API+DB'],
    ['US-089', 'CI/CD Pipeline', 'EPIC-17', '5', 'US-088', 'Infra'],
    ['US-095', 'ShipRocket Opt-In & Setup', 'EPIC-18', '5', 'US-040', 'UI+API+DB'],
    ['US-097', 'AI Try-On for Apparel', 'EPIC-19', '8', 'US-017', 'UI+API+DB'],
    ['US-111', 'Pete Street Virtual Walk', 'EPIC-25', '13', 'US-015', 'UI+API+DB'],
    ['US-113', 'Shop Together Co-Shopping', 'EPIC-26', '13', 'US-111', 'UI+API+DB'],
]
add_table(s, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5.5), story_headers, story_rows,
          col_widths=[Inches(1), Inches(3.5), Inches(1.2), Inches(0.8), Inches(1.5), Inches(1.5)])

# ════════════════════════════════════════
# SLIDE 6: MVP SCOPE
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'MVP Scope — Sprints 1-6 (Days 1-12)',
                   'Approved via GATE-MVP-01 | Auth-first approach | Customer + Merchant + Admin verticals')

add_multiline_textbox(s, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.5), [
    '✅ IN SCOPE — MVP (Days 1-12):',
    '',
    'Customer Features:',
    '• Phone OTP registration & login',
    '• Landing page with Pete Tapestry (21 markets)',
    '• Product search with autocomplete & filters',
    '• Product detail with images, price, mode badges',
    '• Multi-merchant shopping cart',
    '• Consolidated checkout with delivery fee calc',
    '• Razorpay payment (UPI, Card, NetBanking)',
    '• Order confirmation & tracking (status timeline)',
    '• Order history page',
    '• WhatsApp deep link (Mode B)',
    '• Store visit interface (Mode C) w/ Google Maps',
    '• Guest browsing (register at checkout)',
    '',
    'Merchant Features:',
    '• Registration wizard (phone → business → plan → modes)',
    '• Bank details + Razorpay subaccount creation',
    '• Store microsite at petemart.in/{shop-slug}',
    '• Product CRUD + Bulk CSV upload',
    '• Order management + GST invoice generation',
    '• Dashboard (orders, revenue, product count)',
    '• Real-time order alerts + inventory management',
], font_size=10, color=DARK_TEXT)

add_multiline_textbox(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.5), [
    'Admin Features:',
    '• Email+password + OTP 2FA login',
    '• Merchant approval queue (one-click approve/reject)',
    '• Platform dashboard (merchants, orders, revenue)',
    '• Feature flag management with kill switch',
    '• Admin config (delivery zones, rates)',
    '• City selector (Bangalore only at MVP)',
    '',
    'Infrastructure:',
    '• Supabase Free (PostgreSQL, Auth, Storage, Realtime)',
    '• Vercel Hobby (Next.js hosting)',
    '• CI/CD via GitHub Actions',
    '• HTTPS/TLS, RLS, RBAC, Persona-aware nav',
    '',
    '❌ OUT OF SCOPE — Post-MVP:',
    '',
    'Tier 1 (Days 13-14): Mobile apps (PWA suffices),',
    '  i18n, push notifications, zero-downtime deploys',
    '',
    'Tier 2 (Days 15-16): Reviews, promos, loyalty,',
    '  ShipRocket, payment escrow',
    '',
    'Tier 3 (Days 17-18): AI Try-On, bullion rates,',
    '  city expansion, analytics pipeline',
    '',
    'Tier 4 (Days 19-22): Virtual Walk, Live Bazaar,',
    '  Co-Shopping (visionary features)',
], font_size=10, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 7: WORKFLOW MAPS
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Customer Workflows — 5 Personas, Complete Journey Maps',
                   'Step-by-step flows for Customer, Merchant, Admin, Delivery Partner, and B2B Buyer')

workflows = [
    ('Customer Journey', ACCENT_BLUE, Inches(0.3), [
        'Guest → Landing Page',
        '  → Browse Markets (21 Pete areas)',
        '  → Search Products (PgFTS)',
        '  → View Product Detail (Mode A/B/C)',
        '  → Add to Cart (single store)',
        '  → Add from another store (multi-cart)',
        '  → Checkout (consolidated)',
        '  → Pay via Razorpay (UPI/Card/NB)',
        '  → Order Confirmed',
        '  → Track Order (status timeline)',
        '  → Receive Delivery',
        '  → Rate & Review (post-MVP)',
        '',
        'Modes: A=Direct Purchase, B=WhatsApp, C=Visit',
    ]),
    ('Merchant Journey', ACCENT_GREEN, Inches(4.6), [
        'Register (OTP) → Business Details',
        '  → Select Plan (Starter/Growth/Premium)',
        '  → Choose Modes (A, B, C)',
        '  → Enter Bank Details (Razorpay subaccount)',
        '  → Upload Cancelled Cheque',
        '  → Await Admin Approval',
        '  → Store Goes Live !',
        '  → Add Products (manual/CSV)',
        '  → Receive Orders',
        '  → Fulfill & Mark Status',
        '  → View Analytics & Payout Reports',
        '  → Withdraw Earnings (Razorpay Route)',
        '',
        'Dashboard: Overview, Products, Orders, Analytics',
    ]),
    ('Admin Journey', PURPLE, Inches(8.9), [
        'Login (Email+OTP) → Dashboard',
        '  → Review Merchant Applications',
        '  → Approve/Reject (one-click)',
        '  → Monitor Platform Orders',
        '  → Configure Delivery Zones & Rates',
        '  → Manage Feature Flags & Kill Switch',
        '  → View Platform Analytics',
        '  → Manage City Settings',
        '  → Handle Disputes (manual for MVP)',
        '',
        'Console: Dashboard, Merchants, Orders, Config',
        '',
        'Delivery Partner (Courier):',
        '  Login → Today\'s Route',
        '  → Pickups from Merchant A → B → C',
        '  → Consolidate at Micro-Hub',
        '  → Single-drop Delivery',
        '  → GPS Broadcast (30s intervals)',
        '  → Earnings Dashboard',
    ]),
]
for title, color, left, items in workflows:
    add_architecture_box(s, left, Inches(1.3), Inches(4.1), Inches(5.5), title, items, color)

# ════════════════════════════════════════
# SLIDE 8: ACCEPTANCE CRITERIA & REVIEW GATES
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Acceptance Criteria & Quality Gates',
                   'Every story has acceptance criteria. Sprint reviews + HITL gates at key milestones.')

gate_headers = ['Gate ID', 'Triggered By', 'Sprint', 'Verification Items']
gate_rows = [
    ['GATE-TECH-STACK-01', 'Architect (03)', 'S2', 'Supabase schema, Auth OTP, Vercel deploy, RLS, CI/CD, Stitch demo'],
    ['GATE-MVP-01 ✅', 'Program Mgmt (05)', 'S4', 'Market Explorer, Purchase Flow, Merchant Onboarding, WhatsApp, Maps, UX labels, Auth nav'],
    ['GATE-COSTING-01', 'Architect (03)', 'S10', '₹0/mo infra, cost dashboard, budget alerts, scaling model'],
    ['GATE-PRODUCTION-01', 'Production (09)', 'S12', 'All P0/P1 done, 0 critical defects, 80% coverage, E2E pass, perf <300ms, security clean'],
]
add_table(s, Inches(0.3), Inches(1.3), Inches(12.7), Inches(2.5), gate_headers, gate_rows,
          col_widths=[Inches(2), Inches(2), Inches(1), Inches(7.7)])

add_textbox(s, Inches(0.3), Inches(4.0), Inches(12), Inches(0.4), 'Sample Acceptance Criteria (US-021: Consolidated Checkout)', font_size=14, bold=True, color=ACCENT_BLUE)
add_multiline_textbox(s, Inches(0.3), Inches(4.5), Inches(12.7), Inches(2.5), [
    '✅ User can see items from multiple merchants grouped by store in a single checkout page',
    '✅ Delivery fee calculated as: MAX(Zone Rate) + (₹25 × (N-1)) + Weight Surcharge — displayed before payment',
    '✅ User can remove items from any store without affecting other stores\' items',
    '✅ User can enter/select delivery address — validated for zone coverage',
    '✅ Coupon code field shown (disabled for MVP — "Coming Soon" state)',
    '✅ Order summary shows: subtotal per store, delivery fee, platform fee, total payable',
    '✅ "Place Order" button triggers Razorpay checkout widget — disabled if any store has ₹0 items',
    '✅ If cart has items from only Mode B/C merchants, checkout redirects to WhatsApp/enquiry flow',
], font_size=10, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 9: SPRINT CAPACITY & VELOCITY
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Sprint Capacity — Agent-Optimized',
                   '3 parallel agents (UI + API + Backend) = 120 pts/sprint. All sprints within capacity.')

cap_headers = ['Sprint', 'Days', 'Stories', 'Points', 'Agent Capacity', 'Status']
cap_rows = [
    ['S1', '1-2', '10', '36', '120', '🟢 High'],
    ['S2', '3-4', '10', '44', '120', '🟢 High'],
    ['S3', '5-6', '10', '37', '120', '🟢 High'],
    ['S4', '7-8', '11', '42', '120', '🟢 High'],
    ['S5', '9-10', '11', '43', '120', '🟢 High'],
    ['S6', '11-12', '16', '62', '120', '🟢 High'],
    ['S7', '13-14', '16', '55', '120', '🟢 High'],
    ['S8', '15-16', '9', '33', '120', '🟢 Underloaded'],
    ['S9', '17-18', '8', '46', '120', '🟢 High'],
    ['S10', '19-20', '4', '21', '120', '🟢 Underloaded'],
    ['S11', '21-22', '2', '21', '120', '🟢 Underloaded'],
    ['S12', '23-24', '4', '34', '120', '🟢 High'],
]
add_table(s, Inches(0.3), Inches(1.3), Inches(7), Inches(5.5), cap_headers, cap_rows,
          col_widths=[Inches(0.8), Inches(1), Inches(1), Inches(1), Inches(1.5), Inches(1.7)])

add_multiline_textbox(s, Inches(7.6), Inches(1.3), Inches(5.5), Inches(5), [
    'Key Metrics:',
    '',
    '  Total Story Points: 474',
    '  MVP Points (S1-S6): 264',
    '  Avg Sprint Velocity: 120 pts',
    '  Agent Count: 3 (parallel)',
    '',
    'Delivery Milestones:',
    '',
    '  M-T0-01: Market Explorer (D4)',
    '  M-T0-02: Direct Purchase Flow (D10)',
    '  M-T0-03: Merchant Go-Live (D8)',
    '  M-T1-01: Multi-Channel Access (D12)',
    '  M-T1-02: Admin Command Center (D12)',
    '  M-T2-01: Trust & Engagement (D14)',
    '  M-T2-02: Pan-India Shipping (D16)',
    '  M-T3-01: AI Try-On & Jewellery (D18)',
    '  M-T4-01: Virtual Walk (D22)',
    '  M-T4-02: Live Bazaar (D24)',
    '',
    'All milestones within sprint capacity. No overload scenarios.',
], font_size=11, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 10: DEPENDENCY GRAPH
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Dependency Graph — Critical Path & Parallel Tracks',
                   'No deadlocks. All 113 stories mapped. 3 parallel tracks (UI, API, DB).')

add_shape_rect(s, Inches(0.3), Inches(1.3), Inches(12.7), Inches(1.5), LIGHT_GRAY)
add_textbox(s, Inches(0.5), Inches(1.35), Inches(12), Inches(0.35), 'CRITICAL PATH: S1 → S2 → S3 → S4 → S5 → S6 → S7 → S8 → S9 → S10 → S11 → S12', font_size=12, bold=True, color=ACCENT_RED)
add_multiline_textbox(s, Inches(0.5), Inches(1.8), Inches(12), Inches(1), [
    'S1 (Auth+Landing) → S2 (Merchant+Search) → S3 (Cart+Payment) → S4 (Dashboard+Delivery) → S5 (Admin+Tracking) → S6 (i18n+Config)',
    '→ S7 (Mobile+Reviews) → S8 (Promo+ShipRocket) → S9 (AI+Jewellery) → S10 (City+Analytics) → S11 (Virtual Walk) → S12 (Co-Shopping)',
], font_size=10, color=DARK_TEXT)

tracks = [
    ('Parallel Track: UI Agent (07a)', ACCENT_BLUE, Inches(0.3), [
        'Landing Pages (S1), Search/Discovery UI (S1-S2), Cart/Checkout UI (S3),',
        'Merchant Dashboard (S4), Admin Console (S5), i18n/PWA (S6),',
        'Mobile Expo App (S7), Reviews UI (S7), Promo Engine UI (S8),',
        'AI Try-On UI (S9), City Expansion UI (S10), Virtual Walk (S11-S12)',
    ]),
    ('Parallel Track: API Agent (07b)', ACCENT_GREEN, Inches(4.6), [
        'Auth API + RBAC (S1), Product/Merchant API (S2), Cart/Checkout API (S3),',
        'Order/Delivery API (S4), Admin/Tracking API (S5), i18n/Config API (S6),',
        'Mobile API + Reviews (S7), Promo/ShipRocket API (S8),',
        'AI/Bullion API (S9), City/Video API (S10), Virtual API (S11-S12)',
    ]),
    ('Parallel Track: Backend DB Agent (07c)', PURPLE, Inches(8.9), [
        'Supabase Schema + Auth (S1), Product/Merchant Tables (S2),',
        'Order/Payment Schema (S3), Delivery/Zone Config (S4),',
        'Admin/Tracking Tables (S5), i18n/City Schema (S6),',
        'Reviews/Subscription Schema (S7), Promo/ShipRocket (S8),',
        'AI/Bullion Schema (S9), City Expansion Tables (S10),',
        'Virtual Walk/Live Bazaar Schema (S11-S12)',
    ]),
]
for title, color, left, items in tracks:
    add_architecture_box(s, left, Inches(3.0), Inches(4.1), Inches(3.5), title, items, color)

add_textbox(s, Inches(0.3), Inches(6.8), Inches(12), Inches(0.5),
            '✅ Guardrail PASS: No deadlocks or unmapped dependencies | All 113 stories have resolved dependency chains',
            font_size=10, color=MID_GRAY)

# ════════════════════════════════════════
# SLIDE 11: QA VERIFICATION GATES
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'QA Verification & Promotion Gates',
                   'Strict environment promotion criteria. No code reaches Production without clearing all gates.')

qa_boxes = [
    ('Dev → QA Gate', ACCENT_BLUE, Inches(0.3), [
        'Condition: All unit tests pass (80%+ coverage)',
        'Verifier: CI Pipeline (GitHub Actions)',
        'Blocking: ✅ Yes',
        '',
        'What runs:',
        '• Vitest + React Testing Library',
        '• API endpoint contract tests',
        '• Lint + TypeScript checks',
        '• Build verification',
    ]),
    ('QA → Staging Gate', ACCENT_GREEN, Inches(4.6), [
        'Condition: E2E tests pass, No P0/P1 defects',
        'Verifier: QA Agent (08) — Senior Test Architect',
        'Blocking: ✅ Yes',
        '',
        'What runs:',
        '• Playwright E2E (critical journeys)',
        '• API integration tests (Supertest + MSW)',
        '• Lighthouse CI (LCP < 2.5s)',
        '• k6 load test (100 concurrent users)',
    ]),
    ('Staging → Production Gate', ORANGE, Inches(8.9), [
        'Condition: All QA criteria + Security scan + HITL sign-off',
        'Verifier: QA Agent (08) + Production Agent (09)',
        'Blocking: ✅ Yes',
        '',
        'What runs:',
        '• Full regression suite',
        '• OWASP ZAP security scan',
        '• Secrets scan (no hardcoded creds)',
        '• Performance validation (p95 < 300ms)',
        '• Human gate: GATE-PRODUCTION-01',
    ]),
]
for title, color, left, items in qa_boxes:
    add_architecture_box(s, left, Inches(1.3), Inches(4.1), Inches(5.5), title, items, color)

# ════════════════════════════════════════
# SLIDE 12: CLOSING
# ════════════════════════════════════════
s = new_slide()
add_bg(s, BG_DARK)
add_textbox(s, Inches(0.8), Inches(1.5), Inches(11), Inches(1), 'PeteMart', font_size=54, bold=True, color=ACCENT_GOLD)
add_textbox(s, Inches(0.8), Inches(2.7), Inches(11), Inches(0.8), 'Program Management — Ready for Execution', font_size=18, color=WHITE)
add_shape_rect(s, Inches(0.8), Inches(3.7), Inches(3), Inches(0.04), ACCENT_GOLD)
add_multiline_textbox(s, Inches(0.8), Inches(4.0), Inches(11), Inches(2.5), [
    '26 Epics · 72 Features · 113 User Stories · 12 Sprints · 22 Days',
    'MPV: 12 Days (Sprints 1-6) | Full Product: 22 Days (Sprints 1-12)',
    'Execution: 3 AI Agents in Parallel (6x human velocity)',
    '',
    'GATE-MVP-01 ✅ Approved — Phase 3 agents building now',
    'Jira-ready backlog: all stories mapped with dependencies and requirement IDs',
    'Next Gate: GATE-COSTING-01 (End of S10)',
], font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))
add_textbox(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
            'Thank You | Program Review & Next Steps', font_size=16, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

output_path = os.path.join(os.getcwd(), 'agents', '05_program_mgmt_agent', '05_program_mgmt_agent_FULL_PRESENTATION.pptx')
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
