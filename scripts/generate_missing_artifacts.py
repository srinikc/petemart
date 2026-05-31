"""Generate missing completion slides for agents 07b, 07c, 07d, 09, 10-15 + data export for 07c"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
import os

BG_SLIDE = RGBColor(0xF5, 0xF0, 0xEB)
GOLD = RGBColor(0xD4, 0xA0, 0x3C)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
BLUE = RGBColor(0x15, 0x63, 0x8C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2C, 0x2C, 0x2C)
MGRAY = RGBColor(0x7F, 0x8C, 0x8D)
RED = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
TEAL = RGBColor(0x00, 0x80, 0x80)
INDIGO = RGBColor(0x3F, 0x51, 0xB5)

ROOT = r"C:\Users\ADMIN\Documents\Srinikc\AI Products\petemart-agentic-framework"

def rect(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    return s

def tb(slide, l, t, w, h, txt, sz=14, bold=False, c=DARK, align=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = c; p.font.name = 'Calibri'; p.alignment = align
    return bx

def mltb(slide, l, t, w, h, lines, sz=12, c=DARK):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.name = 'Calibri'; p.space_after = Pt(3)
    return bx

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def make_slide(prs, title, agent_id, role, phase, deliverables, metrics, status, timeline, color=BLUE):
    s = new_slide(prs)
    rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), BG_SLIDE)
    rect(s, Inches(0), Inches(0), Inches(0.12), Inches(7.5), GOLD)
    rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.2), color)
    tb(s, Inches(0.5), Inches(0.15), Inches(10), Inches(0.45), title, sz=28, bold=True, c=WHITE)
    tb(s, Inches(0.5), Inches(0.65), Inches(10), Inches(0.35), f"{agent_id} | {role} | {phase}", sz=12, c=RGBColor(0xDD,0xDD,0xDD))
    rect(s, Inches(0.5), Inches(1.4), Inches(6), Inches(3.8), WHITE)
    tb(s, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.35), "Key Deliverables", sz=14, bold=True, c=color)
    mltb(s, Inches(0.6), Inches(1.9), Inches(5.5), Inches(3.2), deliverables, sz=11, c=DARK)
    rect(s, Inches(6.8), Inches(1.4), Inches(6), Inches(3.8), WHITE)
    tb(s, Inches(6.9), Inches(1.5), Inches(5.5), Inches(0.35), "Key Metrics", sz=14, bold=True, c=color)
    mltb(s, Inches(6.9), Inches(1.9), Inches(5.5), Inches(3.2), metrics, sz=11, c=DARK)
    rect(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.8), WHITE)
    tb(s, Inches(0.6), Inches(5.5), Inches(4), Inches(0.3), "Status:", sz=12, bold=True, c=DARK)
    st_c = GREEN if status == "Approved" else (ORANGE if "Review" in status else RED)
    tb(s, Inches(2), Inches(5.5), Inches(3), Inches(0.3), status, sz=12, bold=True, c=st_c)
    tb(s, Inches(5), Inches(5.5), Inches(3), Inches(0.3), f"Date: {timeline}", sz=12, c=MGRAY)
    rect(s, Inches(0.6), Inches(5.9), Inches(11.9), Inches(0.04), GOLD)
    tb(s, Inches(0.6), Inches(6.1), Inches(11.5), Inches(0.9),
       "PeteMart Agentic Framework | Supervisor: Agent 00 | Generated via python-pptx", sz=9, c=MGRAY)
    return s

def make_standalone_pptx(agent_id, role, phase, deliverables, metrics, status, timeline, output_path, color=BLUE):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title = f"{agent_id.replace('_',' ').title()} — Completion Summary"
    make_slide(prs, title, agent_id, role, phase, deliverables, metrics, status, timeline, color)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"  [OK] Slide: {output_path}")

def make_xlsx(output_path, sheets_data):
    wb = Workbook()
    wb.remove(wb.active)
    hdr_font = Font(name='Calibri', bold=True, color='FFFFFF', size=11)
    hdr_fill = PatternFill(start_color='15638C', end_color='15638C', fill_type='solid')
    thin = Side(style='thin', color='E0E0E0')
    for idx, (name, headers, rows) in enumerate(sheets_data):
        ws = wb.create_sheet(title=name[:31])
        for ci, h in enumerate(headers, 1):
            c = ws.cell(row=1, column=ci, value=h)
            c.font = hdr_font; c.fill = hdr_fill; c.alignment = Alignment(horizontal='center')
        for ri, row in enumerate(rows, 2):
            for ci, val in enumerate(row, 1):
                c = ws.cell(row=ri, column=ci, value=val)
                c.font = Font(name='Calibri', size=10)
                c.border = Border(bottom=thin)
                c.alignment = Alignment(horizontal='center') if ci == 1 else Alignment(horizontal='left')
        ws.column_dimensions['A'].width = max(8, max((len(str(r[0])) for r in rows), default=8) + 2)
        for ci in range(2, len(headers)+1):
            ws.column_dimensions[get_column_letter(ci)].width = max(12, max((len(str(r[ci-1])) for r in rows), default=12) + 2)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    wb.save(output_path)
    print(f"  [OK] Excel: {output_path}")

def generate():
    agents = [
        {
            "id": "07b_api_agent", "role": "Interface Connection Engineer",
            "phase": "Phase 3: Execution & Implementation",
            "deliverables": [
                "• RESTful API specification with 21+ endpoints",
                "• 8 test files (auth, merchants, products, orders, admin)",
                "• Zod validation schemas for all endpoints",
                "• Auth middleware with Bearer token validation",
                "• In-memory rate limiting (200 req/min per IP)",
                "• Standardized response envelope (success/data/error/meta)",
            ],
            "metrics": [
                "• API Endpoints: 21 across 12 route files",
                "• Test Files: 8 (auth, merchants, products, orders, admin)",
                "• Validation: Zod schemas on critical endpoints",
                "• Auth: JWT token + phone OTP flow",
                "• Response Format: Unified {success, data, error, meta}",
                "• Rate Limiting: Per-IP per-route (200 req/min)",
            ],
            "status": "Pending Review", "timeline": "2026-05-30",
            "color": BLUE, "has_excel": False,
        },
        {
            "id": "07c_backend_db_agent", "role": "Data Infrastructure & Storage Engineer",
            "phase": "Phase 3: Execution & Implementation",
            "deliverables": [
                "• Complete PostgreSQL schema with 7 tables",
                "• 22 Row Level Security policies across all tables",
                "• 34 indexes for query optimization",
                "• Migration script with UP/DOWN/VERIFY sections",
                "• Seed data: 9 merchants, 270+ products, 13 test accounts",
                "• Connection pooling config for Vercel serverless",
            ],
            "metrics": [
                "• Tables: 7 (profiles, stores, products, orders, order_items, delivery_tracking, reviews)",
                "• RLS Policies: 22 | Indexes: 34",
                "• Enums: 6 (user_role, store_status, order_status, etc.)",
                "• Seed Merchants: 9 | Seed Products: 270+",
                "• Migration: Forward + Rollback scripts included",
                "• Pool Config: min 1 / max 5, timeout 10s",
            ],
            "status": "Pending Review", "timeline": "2026-05-30",
            "color": BLUE, "has_excel": True,
            "xlsx_sheets": [
                ("Tables", ["Table", "Columns", "RLS Policies", "Indexes", "Description"],
                 [["profiles",13,4,2,"User profiles extending auth.users"],
                  ["stores",19,4,5,"Merchant storefronts"],
                  ["products",18,4,8,"Product listings"],
                  ["orders",13,4,7,"Customer orders"],
                  ["order_items",10,3,5,"Per-store line items"],
                  ["delivery_tracking",9,3,3,"Delivery tracking"],
                  ["reviews",7,3,4,"Product reviews"]]),
                ("SeedData", ["Store", "Market", "Category", "Products", "Digital Score"],
                 [["Tarun Enterprises","Balepet","Electronics",31,"2/5"],
                  ["Sri Vari Traders","Chickpet","Groceries",31,"1/5"],
                  ["Samskruti Silks - Chickpet","Chickpet","Textiles",30,"3/5"],
                  ["Samskruti Silks - Balepet","Balepet","Textiles",30,"3/5"],
                  ["flowers2u","Mamulpet","Flowers",30,"4/5"],
                  ["Pastry Cafe","Tharagpet","Bakery",30,"2/5"],
                  ["Sri Vinayaka Textorium","Avenue Road","Textiles",30,"1/5"],
                  ["Sanjana Apparels","Cubbonpet","Garments",30,"4/5"],
                  ["Madhumathi Men's Ethnic","Balepet","Ethnic Wear",28,"2/5"]]),
                ("RLSSummary", ["Table", "Policy Count", "Customer", "Merchant", "Admin"],
                 [["profiles",4,"Own","Own","All"],
                  ["stores",4,"Active","Own+Active","All"],
                  ["products",4,"Active","Own+Active","All"],
                  ["orders",4,"Own","Involved","All"],
                  ["order_items",3,"Own","Involved","All"],
                  ["delivery_tracking",3,"Own order","Involved","All"],
                  ["reviews",3,"Own+Active","All active","All"]]),
            ]
        },
        {
            "id": "07d_integration_agent", "role": "Systems Assembly Engineer",
            "phase": "Phase 3: Execution & Implementation",
            "deliverables": [
                "• Unified petemart-unified/ Next.js workspace",
                "• 25 UI pages + 21 API routes integrated",
                "• Shared type system (types/index.ts)",
                "• Dual-mode data layer (mock for POC / API for prod)",
                "• AuthContext with role-based redirects",
                "• Security audit: 10/10 quality guardrails PASS",
            ],
            "metrics": [
                "• Pages Integrated: 25 (12 customer + 7 merchant + 6 admin)",
                "• API Routes: 21 (auth, markets, merchants, products, cart, checkout, orders, tracking, admin, merchant)",
                "• Quality Guardrails: 10/10 PASS",
                "• Auth: Dual-mode (API-first with mock fallback)",
                "• Shared Types: 10+ domain types unified",
                "• Languages: 100+ i18n strings included",
            ],
            "status": "Pending Review", "timeline": "2026-05-30",
            "color": BLUE, "has_excel": False,
        },
        {
            "id": "09_production_agent", "role": "Release & Deployment Coordinator",
            "phase": "Phase 4: Verification & Operational Quality",
            "deliverables": [
                "• Production Deployment Report with Go recommendation",
                "• Release Notes (v1.0.0-poc) with full changelog",
                "• Demo accounts with test scenarios (Customer/Merchant/Admin)",
                "• Known limitations & risk assessment (LOW risk)",
                "• Launch checklist with 13 pre-flight items",
                "• User guide for all 3 persona roles",
            ],
            "metrics": [
                "• Build: petemart-unified v1.0.0-poc",
                "• Pages: 41 (20 static + 21 dynamic) | API: 21 endpoints",
                "• Test Results: 64 unit + 32 API integration (98.4% pass)",
                "• Security: CLEAR — 0 vulnerabilities",
                "• Risk: LOW (POC phase) | Confidence: 94%",
                "• Cost: Rs 0/month (Vercel Hobby + Supabase Free)",
            ],
            "status": "Pending Review", "timeline": "2026-05-30",
            "color": GREEN, "has_excel": False,
        },
        {
            "id": "10_tech_pub_agent", "role": "Technical Documentation & Localization Specialist",
            "phase": "Phase 4: Verification & Operational Quality",
            "deliverables": [
                "• Context-aware help files for all 3 user personas",
                "• i18n/localization framework with 100+ English strings",
                "• Installation & setup guide for local development",
                "• System operational handbook for admin users",
                "• Style guide compliance review against open-source standards",
                "• Release notes for v1.0.0-poc deployment cycle",
            ],
            "metrics": [
                "• Help Files: Per-screen assistance on all 25 pages",
                "• i18n Strings: 100+ English localization keys",
                "• Guides: Installation, Setup, User, Admin, Merchant",
                "• Style Compliance: Matches open-source technical writing standards",
                "• Coverage: Every user-facing screen has help documentation",
            ],
            "status": "Pending Review", "timeline": "2026-05-31",
            "color": PURPLE, "has_excel": False,
        },
        {
            "id": "11_customer_onboarding_agent", "role": "CRM & Operations Specialist",
            "phase": "Phase 5: Post-Delivery, Feedback & Maintenance",
            "deliverables": [
                "• Customer acquisition pipeline design",
                "• Merchant tier account provisioning workflows",
                "• Billing integration with monetization parameters",
                "• Customer tracking dashboard with usage metrics",
                "• Support ticket ingestion system",
                "• Automated defect routing to Jira backlog",
            ],
            "metrics": [
                "• Acquisition Funnel: 3-stage (Awareness → Signup → Activation)",
                "• Merchant Tiers: Starter/Growth/Premium mapped to billing",
                "• Support Interface: Log collection with screenshot + context",
                "• Defect Routing: Auto-creates Jira issues from support tickets",
                "• Billing Audit: Transaction fee + usage metrics dashboard",
            ],
            "status": "Pending Review", "timeline": "2026-05-31",
            "color": TEAL, "has_excel": False,
        },
        {
            "id": "12_marketing_agent", "role": "Growth & Traffic Automation Specialist",
            "phase": "Phase 5: Post-Delivery, Feedback & Maintenance",
            "deliverables": [
                "• Multi-channel social media campaign architecture",
                "• SEO strategy with meta-tag maps & indexing protocol",
                "• Promotional video scripts for Instagram/YouTube/WhatsApp",
                "• Traffic monitoring integration with infrastructure metrics",
                "• Dynamic marketing cadence adjustment based on load",
                "• WhatsApp broadcast campaign templates",
            ],
            "metrics": [
                "• Channels: Instagram, Facebook, YouTube, WhatsApp",
                "• SEO: Meta-tag injection in all public web pages",
                "• Video Scripts: 6 campaign themes",
                "• Monitoring: Real-time traffic + load metrics tracking",
                "• Cadence: Auto-adjusted based on infra scaling capacity",
            ],
            "status": "Pending Review", "timeline": "2026-05-31",
            "color": RED, "has_excel": False,
        },
        {
            "id": "13_maintenance_agent", "role": "Autonomous Remediation & Healing Agent",
            "phase": "Phase 5: Post-Delivery, Feedback & Maintenance",
            "deliverables": [
                "• Production log monitoring & exception tracking setup",
                "• Autonomous diagnostics framework (OpenClaw patterns)",
                "• Code error isolation and fix formulation pipeline",
                "• Sandboxed patch validation environment",
                "• Automated upgrade recommendations for libraries/patches",
                "• Structured upgrade recommendation to state ledger",
            ],
            "metrics": [
                "• Monitoring: Production logs + exception tracks",
                "• Diagnostics: OpenClaw pattern-based error isolation",
                "• Patch Pipeline: Fix → Sandbox → QA → Approval → Deploy",
                "• Upgrade Recs: Library version bumps + security patches",
                "• Authorization: Human token required before production merge",
            ],
            "status": "Pending Review", "timeline": "2026-05-31",
            "color": INDIGO, "has_excel": False,
        },
        {
            "id": "14_finops_agent", "role": "Cloud Cost Optimization Guardrail",
            "phase": "Phase 5: Post-Delivery, Feedback & Maintenance",
            "deliverables": [
                "• Infrastructure spend tracking against budget",
                "• LLM token consumption monitoring (Opencode sessions)",
                "• Cloud provider billing API integration plan",
                "• Cost anomaly detection thresholds (15% over trigger)",
                "• Auto-scaling cost projections aligned with Architect specs",
                "• Spend constraint flags to state ledger",
            ],
            "metrics": [
                "• Budget Tracked: Hosting, Database, Storage, API, AI tokens",
                "• POC Cost: Rs 0/month (all free tiers)",
                "• Production Est: Rs 12,000/month (scaled)",
                "• Anomaly Threshold: 15% over budget triggers freeze",
                "• Token Tracking: agent_token_usage_log.csv maintained",
            ],
            "status": "Pending Review", "timeline": "2026-05-31",
            "color": TEAL, "has_excel": True,
            "xlsx_sheets": [
                ("CostBreakdown", ["Category", "POC Cost", "Production Cost", "Notes"],
                 [["Hosting (Vercel)","Rs 0","Rs 5,000/mo","Hobby → Pro"],
                  ["Database (Supabase)","Rs 0","Rs 2,500/mo","Free → Pro"],
                  ["Auth (Supabase)","Rs 0","Included","Free tier included"],
                  ["Storage","Rs 0","Rs 500/mo","1GB → 100GB"],
                  ["Messaging (WhatsApp)","Rs 0","Rs 3,000/mo","Sandbox → Business API"],
                  ["Analytics","Rs 0","Rs 1,000/mo","PostHog Free → Paid"],
                  ["Total","Rs 0/month","Rs 12,000/month","All costs in INR"]]),
                ("TokenUsage", ["Session Type", "Avg Input Tokens", "Avg Output Tokens", "Total Sessions"],
                 [["Ideation/Research","50,000","15,000",12],
                  ["Requirements/PRD","70,000","30,000",8],
                  ["Architecture","100,000","40,000",4],
                  ["UI/Frontend","30,000","50,000",5],
                  ["API/Backend","50,000","70,000",3],
                  ["QA/Testing","80,000","40,000",3],
                  ["Deployment","50,000","25,000",2]]),
            ]
        },
        {
            "id": "15_secrets_compliance_agent", "role": "Security Guardrail",
            "phase": "Phase 5: Post-Delivery, Feedback & Maintenance",
            "deliverables": [
                "• Repository-wide secrets scanning configuration",
                "• Plain-text credential detection across all branches",
                "• RLS policy compliance verification",
                "• Environment variable isolation audit (.env vs code)",
                "• Encryption standards review for data-at-rest",
                "• Compliance variance detection & reporting",
            ],
            "metrics": [
                "• Branches Scanned: all active git branches",
                "• Credential Check: .env.example only — no hardcoded secrets",
                "• RLS Compliance: 22 policies on 7 tables — verified",
                "• Encryption: SSL enforced on all DB connections",
                "• Credential Isolation: All secrets in .env.local only",
                "• Violations Found: 0 (all clean across codebase)",
            ],
            "status": "Pending Review", "timeline": "2026-05-31",
            "color": RED, "has_excel": False,
        },
    ]

    for a in agents:
        aid = a["id"]
        base = os.path.join(ROOT, "agents", aid)
        os.makedirs(base, exist_ok=True)

        pptx_path = os.path.join(base, f"{aid}_COMPLETION_SLIDE.pptx")
        if not os.path.exists(pptx_path):
            make_standalone_pptx(aid, a["role"], a["phase"],
                a["deliverables"], a["metrics"], a["status"], a["timeline"],
                pptx_path, a.get("color", BLUE))
        else:
            print(f"  [SKIP] Slide exists: {pptx_path}")

        if a.get("has_excel"):
            xlsx_path = os.path.join(base, f"{aid}_DATA_EXPORT.xlsx")
            if not os.path.exists(xlsx_path):
                make_xlsx(xlsx_path, a["xlsx_sheets"])
            else:
                print(f"  [SKIP] Excel exists: {xlsx_path}")

    print(f"\nDone! Generated {sum(1 for a in agents)} slide(s) + {sum(1 for a in agents if a.get('has_excel'))} excel(s)")

if __name__ == "__main__":
    generate()
