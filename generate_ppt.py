from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import json, os

# ── Brand Colors ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_SLIDE = RGBColor(0xF5, 0xF0, 0xEB)
ACCENT_GOLD = RGBColor(0xD4, 0xA0, 0x3C)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
ACCENT_BLUE = RGBColor(0x15, 0x63, 0x8C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
LIGHT_GRAY = RGBColor(0xF0, 0xEE, 0xE9)
MID_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
TEAL = RGBColor(0x00, 0x7B, 0x7F)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper Functions ──
def add_bg(slide, color=BG_SLIDE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_accent_bar(slide, left=Inches(0), top=Inches(0), width=Inches(0.15), height=Inches(7.5), color=ACCENT_GOLD):
    return add_shape_rect(slide, left, top, width, height, color)

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tb

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=12, color=DARK_TEXT, line_spacing=1.3, bold_first=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        if bold_first and i == 0:
            p.font.bold = True
            p.font.size = Pt(font_size + 2)
        p.space_after = Pt(4)
    return tb

def add_kpi_box(slide, left, top, width, height, number, label, bg_color=WHITE, num_color=ACCENT_GOLD):
    box = add_shape_rect(slide, left, top, width, height, bg_color)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.6),
                str(number), font_size=28, bold=True, color=num_color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.65), width - Inches(0.3), Inches(0.4),
                label, font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    return box

def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE
            p.font.name = 'Calibri'
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = DARK_TEXT
                p.font.name = 'Calibri'
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 0 else LIGHT_GRAY
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    return tbl

def add_chart_bar(slide, left, top, width, height, categories, values, series_name='', color=ACCENT_BLUE):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data).chart
    chart.has_legend = False
    plot = chart.plots[0]
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    return chart

def new_slide():
    slide_layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(slide_layout)

# ════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════
s = new_slide()
add_bg(s, BG_DARK)
add_textbox(s, Inches(0.8), Inches(0.6), Inches(11), Inches(1),
            'PeteMart', font_size=54, bold=True, color=ACCENT_GOLD)
add_textbox(s, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
            'Hyper-Local E-Commerce Marketplace for Old Bangalore\'s Historic Pete Markets', font_size=22, color=WHITE)
add_shape_rect(s, Inches(0.8), Inches(2.5), Inches(3), Inches(0.04), ACCENT_GOLD)
add_multiline_textbox(s, Inches(0.8), Inches(2.8), Inches(11), Inches(2.5), [
    'Management Presentation',
    'Phase 1: Ideation & Requirements — Review Package',
    '',
    'Presented for: Higher Management Review',
    'Date: May 2026',
    'Status: Market Research Complete | 103 Requirements Defined',
], font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

# ════════════════════════════════════════
# SLIDE 2: EXECUTIVE SNAPSHOT
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_accent_bar(s)
add_textbox(s, Inches(0.5), Inches(0.3), Inches(6), Inches(0.7),
            'Executive Snapshot', font_size=32, bold=True, color=ACCENT_BLUE)

kpis = [
    ('5,000+', 'Target Merchants', Inches(0.5)),
    ('21', 'Pete Markets Mapped', Inches(3.7)),
    ('406', 'Profiled Merchants', Inches(6.9)),
    ('103', 'Requirements Defined', Inches(10.1)),
]
for num, label, left in kpis:
    add_kpi_box(s, left, Inches(1.2), Inches(2.8), Inches(1.1), num, label, num_color=ACCENT_GOLD)

add_textbox(s, Inches(0.5), Inches(2.7), Inches(12), Inches(0.5),
            'Platform Overview', font_size=20, bold=True, color=ACCENT_BLUE)

add_multiline_textbox(s, Inches(0.5), Inches(3.3), Inches(6), Inches(3.5), [
    'PeteMart is a multi-tenant digital commerce marketplace migrating 5,000+',
    'traditional physical merchants from 21 historic Pete markets of Old Bangalore',
    'into a unified e-commerce ecosystem (Web + Android + iOS).',
    '',
    'Three Interaction Modes available to every merchant:',
    '  Mode A - Direct Purchase (Cart & Checkout)',
    '  Mode B - WhatsApp Enquiry (Bulk Quotes & Negotiation)',
    '  Mode C - Visit Store (Maps & Offline Footfall)',
    '',
    'Commission: 4% B2C | 1.5% B2B (capped at Rs 500)',
    'Subscriptions: Starter Rs 499 | Growth Rs 999 | Premium Rs 2,499 /mo',
], font_size=12, color=DARK_TEXT)

add_multiline_textbox(s, Inches(6.8), Inches(3.3), Inches(5.5), Inches(3.5), [
    'Key Differentiators vs Competitors:',
    '',
    '• Zero Commission on Mode B & C discovery',
    '• Branded Store Microsite (petemart.in/shop-name)',
    '• Multi-Store Order Consolidation',
    '• AI-Generated Product Reels & Videos',
    '• Hyperlocal Zone-Based Delivery',
    '• WhatsApp Deep-Link Checkouts',
    '',
    'One-Time Setup: Rs 86,000 - Rs 1,50,000',
    'Monthly Burn: Rs 43,000 - Rs 63,000',
    'Break-Even: 150 paid sellers (Month 6)',
], font_size=12, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 3: MARKET COVERAGE
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_accent_bar(s)
add_textbox(s, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            'Old Bangalore Pete Markets — Coverage Map', font_size=32, bold=True, color=ACCENT_BLUE)

markets_data = [
    ['Chickpet', 'Textiles, Silk Sarees', '118', 'Mode A/B/C'],
    ['K R Market', 'Flowers, Fresh Produce', '39', 'Mode A/B/C'],
    ['Sultanpet', 'Paper, Wedding Cards', '42', 'Mode A/B/C'],
    ['Santhusapet', 'Cosmetics, Salon Supplies', '48', 'Mode A/B/C'],
    ['Huriopet', 'Cords, Ropes, Twines', '85', 'Mode A/B/C'],
    ['Balepet', 'Plastics, Kitchenware', '34', 'Mode A/B/C'],
    ['Tharagpet', 'Grains, Pulses, Spices', '31', 'Mode A/B/C'],
    ['BVK Iyengar Rd', 'Electrical Accessories', '29', 'Mode A/B/C'],
    ['Cottonpet', 'Footwear, Garments', '23', 'Mode A/B/C'],
    ['Basettyetpet', 'Lighting, Chandeliers', '23', 'Mode A/B/C'],
    ['Avenue Road', 'Books, Stationery', '18', 'Mode A/B/C'],
    ['Akkipete', 'Pharmaceuticals', '13', 'Mode A/B/C'],
]

headers = ['Pete Market', 'Specialization', 'Mapped Stores', 'Available Modes']
add_table(s, Inches(0.5), Inches(1.2), Inches(7), Inches(5.5), headers, markets_data)

add_multiline_textbox(s, Inches(7.8), Inches(1.2), Inches(5), Inches(4), [
    'Total: 21 Pete Areas Mapped',
    '',
    'Est. 5,000+ Active Merchants Ecosystem',
    '',
    'Categories Covered (12):',
    '  Textiles & Apparel',
    '  Jewellery & Accessories',
    '  Electronics & Electricals',
    '  Books & Stationery',
    '  Grocery, FMCG & Provisions',
    '  Plastics & Household',
    '  Hardware & Construction',
    '  Beauty & Cosmetics',
    '  Pharmaceuticals',
    '  Florist & Fresh Flowers',
    '  Bakery & Cafe',
    '  Sports & Outdoor',
    '',
    'Avg Merchant Rating: 4.53/5.0',
], font_size=11, color=DARK_TEXT)

add_textbox(s, Inches(7.8), Inches(5.5), Inches(5), Inches(0.5),
            'Schema v3.0: 17 enhanced fields per merchant (business hours, payment methods, GST, pricing tiers, facilities, etc.)',
            font_size=10, color=MID_GRAY)

# ════════════════════════════════════════
# SLIDE 4: THREE INTERACTION MODES
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_accent_bar(s)
add_textbox(s, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            'Three Merchant Interaction Modes', font_size=32, bold=True, color=ACCENT_BLUE)
add_textbox(s, Inches(0.5), Inches(1.0), Inches(12), Inches(0.5),
            'Every merchant gets all three modes. Opt-in during onboarding — activate one, two, or all three based on business readiness.',
            font_size=13, color=MID_GRAY)

modes = [
    ('Mode A: Direct Purchase', ACCENT_GREEN, Inches(0.5),
     ['In-app cart & checkout via Razorpay',
      'Order tracking & delivery management',
      'Best for: FMCG, Stationery, Novelties, Bags',
      '4% B2C | 1.5% B2B commission']),
    ('Mode B: WhatsApp Enquiry', ACCENT_BLUE, Inches(4.6),
     ['Trackable WhatsApp deep-link button',
      'Bulk quote requests & price negotiation',
      'Best for: Textiles, Sarees, Hardware,',
      '  Wedding Cards, Electrical Components']),
    ('Mode C: Visit Store', ORANGE, Inches(8.7),
     ['Google Maps directions & store facade',
      'Virtual window display & footfall tracking',
      'Best for: Gold/Silver Jewellery, Premium',
      '  Silks, High-Value Trust Items']),
]

for title, color, left, bullets in modes:
    box = add_shape_rect(s, left, Inches(1.6), Inches(3.8), Inches(4.2), WHITE)
    add_shape_rect(s, left, Inches(1.6), Inches(3.8), Inches(0.55), color)
    add_textbox(s, left + Inches(0.2), Inches(1.65), Inches(3.4), Inches(0.45),
                title, font_size=14, bold=True, color=WHITE)
    add_multiline_textbox(s, left + Inches(0.2), Inches(2.4), Inches(3.4), Inches(3), bullets, font_size=11, color=DARK_TEXT)

add_shape_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.2), LIGHT_GRAY)
add_multiline_textbox(s, Inches(0.7), Inches(6.1), Inches(11.8), Inches(1), [
    'Merchant Onboarding Flow: Sign Up -> Select Modes -> Configure Preferences -> Go Live',
    'Merchants can change modes anytime via dashboard. No lock-in. Zero commission on pure discovery (Modes B & C).',
], font_size=11, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 5: COMPETITIVE UVP
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_accent_bar(s)
add_textbox(s, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            'PeteMart Moat — Competitive Differentiation', font_size=32, bold=True, color=ACCENT_BLUE)

uvp_headers = ['Dimension', 'Amazon / Flipkart', 'JustDial / Indiamart', 'PeteMart (Our Advantage)']
uvp_rows = [
    ['Commission', '15-40%', 'Lead-based fees', '0% on Modes B/C; 1.5-4% on Mode A'],
    ['Store Identity', 'Hidden behind search', 'Static listing', 'Branded microsite (petemart.in/shop)'],
    ['Product Catalog', 'Required & controlled', 'Not available', 'Flexible catalog with AI reels'],
    ['B2B Wholesale Tools', 'Not designed for B2B', 'Bulk enquiry only', 'MOQ toggles, tiered pricing, WhatsApp negotiation'],
    ['Offline Integration', 'Not supported', 'Basic listing', 'Google Maps + Store facade + Virtual Windows'],
    ['Delivery', 'Pan-India logistics', 'No delivery support', 'Zone-based hyperlocal + multi-store consolidation'],
    ['Multi-Store Cart', 'Single seller only', 'N/A', 'Cross-Pete cart (e.g. Silk + Jewellery + Spices)'],
]
add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(3.8), uvp_headers, uvp_rows,
          col_widths=[Inches(1.5), Inches(2.8), Inches(2.8), Inches(5.2)])

add_multiline_textbox(s, Inches(0.5), Inches(5.3), Inches(12), Inches(1.5), [
    'Key PeteMart Differentiators:',
    '• Zero Commission Option: Mode B (WhatsApp) and Mode C (Visit Store) are free — merchants pay only if they want Mode A',
    '• Branded Store Microsite: Every store gets petemart.in/shop-name with QR code, Google indexation, and printable marketing collateral',
    '• Hyperlocal Consolidation: Buy across 21 markets in one cart — consolidation surcharge of only Rs 25 per additional store',
    '• AI-Powered Visual Content: Auto-generated product reels for every new listing — social media ready assets',
], font_size=11, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 6: SUBSCRIPTION PLANS
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_accent_bar(s)
add_textbox(s, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            'Merchant Subscription Plans', font_size=32, bold=True, color=ACCENT_BLUE)

plans = [
    ('STARTER', 'Rs 499/mo', 'Rs 4,990/yr', RGBColor(0x7F, 0x8C, 0x8D), Inches(0.5), [
        '50 Products max',
        'Branded Store Page + QR',
        'All 3 Modes available',
        'Basic analytics',
        'Target: Pure B2B wholesalers',
        'starting digital presence',
        '',
        'Tier Distribution:',
        '80 merchants (19.7%)',
    ]),
    ('GROWTH', 'Rs 999/mo', 'Rs 9,990/yr', ACCENT_BLUE, Inches(4.6), [
        '500 Products max',
        'Cart & Checkout enabled',
        'AI Reels (1/mo) + Analytics',
        'WhatsApp Broadcast (1K)',
        'Razorpay integration',
        'Target: Standard retailers',
        '& wholesalers scaling online',
        '',
        'Tier Distribution:',
        '306 merchants (75.4%)',
    ]),
    ('PREMIUM', 'Rs 2,499/mo', 'Rs 24,990/yr', ACCENT_GOLD, Inches(8.7), [
        'Unlimited Products',
        'Concierge Reel Shoots (4/mo)',
        'Featured Seller Badge',
        'Homepage Banner (2d/mo)',
        'Dedicated Account Manager',
        'Unlimited WhatsApp Broadcast',
        '',
        'Tier Distribution:',
        '20 merchants (4.9%)',
    ]),
]

for title, monthly, annual, color, left, bullets in plans:
    box = add_shape_rect(s, left, Inches(1.4), Inches(3.8), Inches(5.2), WHITE)
    add_shape_rect(s, left, Inches(1.4), Inches(3.8), Inches(0.6), color)
    add_textbox(s, left + Inches(0.2), Inches(1.45), Inches(3.4), Inches(0.45),
                title, font_size=16, bold=True, color=WHITE)
    add_textbox(s, left + Inches(0.2), Inches(2.2), Inches(3.4), Inches(0.4),
                monthly, font_size=24, bold=True, color=color)
    add_textbox(s, left + Inches(0.2), Inches(2.6), Inches(3.4), Inches(0.3),
                annual, font_size=10, color=MID_GRAY)
    add_multiline_textbox(s, left + Inches(0.2), Inches(3.1), Inches(3.4), Inches(3.2), bullets, font_size=11, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 7: REVENUE MODEL
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_accent_bar(s)
add_textbox(s, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            'Revenue Model & Monetization', font_size=32, bold=True, color=ACCENT_BLUE)

rev_headers = ['Revenue Stream', 'Rate', 'Notes']
rev_rows = [
    ['B2C Commission (Mode A)', '4.0% per order', 'Standard retail checkout transactions'],
    ['B2B Wholesale Commission', '1.5% (capped at Rs 500)', 'Bulk wholesale orders via Mode A'],
    ['Payment Gateway Fee', '2.0% pass-through', 'Razorpay processing (pass-through)'],
    ['Catalog Digitization', 'Rs 10/product', 'Done-For-You product upload service'],
    ['AI Reel Creation', 'Rs 999/video', 'Concierge reel shoots for premium content'],
    ['Sponsored CPC Listings', 'Rs 2.00/click', 'Search result boost for featured products'],
    ['Homepage Banner Ads', 'Rs 499/day', 'Premium homepage real estate'],
    ['Courier Platform Share', '15% of delivery fee', 'Zone-based delivery fee split (85/15)'],
]
add_table(s, Inches(0.5), Inches(1.2), Inches(7.5), Inches(3.5), rev_headers, rev_rows)

add_multiline_textbox(s, Inches(8.3), Inches(1.2), Inches(4.5), Inches(3.5), [
    'Projected Revenue Mix (at 400 sellers):',
    '',
    'Subscription MRR:   Rs 3,75,000',
    'Transaction Fees:   Rs 1,20,000',
    'Value-Added Svc:    Rs    45,000',
    'Delivery Revenue:   Rs    30,000',
    '─────────────────────────────',
    'Total Est. MRR:     Rs 5,70,000',
    '',
    '(Year 1 Target: 400 paid sellers)',
    '',
    'Year 2 Target:',
    '2,000 sellers → Rs 17L/mo MRR',
    '→ Rs 2 Crore ARR',
], font_size=11, color=DARK_TEXT)

add_textbox(s, Inches(0.5), Inches(5.0), Inches(12), Inches(0.5),
            'Delivery Cost Model: Zone 1 (0-3km): Rs 40 retail / Rs 60 wholesale | Zone 2 (3-7km): Rs 70 / Rs 110 | Zone 3 (7+km): Rs 110 / Rs 180',
            font_size=10, color=MID_GRAY)

# ════════════════════════════════════════
# SLIDE 8: FINANCIAL PROJECTIONS
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_accent_bar(s)
add_textbox(s, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            'Financial Projections & Break-Even Analysis', font_size=32, bold=True, color=ACCENT_BLUE)

# Setup Costs
setup_headers = ['Category', 'Cost (Rs)']
setup_rows = [
    ['Company Registration & Trademarks', '12,500 - 24,000'],
    ['App Store Listings (Apple + Google)', '10,400'],
    ['Legal Documentation', '25,000 - 40,000'],
    ['Cloud, Supabase & DB Provisioning', '10,000 - 20,000'],
    ['First 50 Onboarding Visits', '25,000 - 50,000'],
    ['Total Setup Cost', '86,000 - 1,50,000'],
]
add_table(s, Inches(0.5), Inches(1.2), Inches(5.5), Inches(2.8), setup_headers, setup_rows)

# Monthly Burn
burn_headers = ['Category', 'Monthly (Rs)']
burn_rows = [
    ['Cloud Infra (Supabase, Railway, Vercel)', '1,300 - 4,000'],
    ['Communications (OTP, SMS, WhatsApp)', '3,000 - 5,000'],
    ['Google Maps API', '1,500 - 3,000'],
    ['AI API Suite (Claude, OpenAI, etc.)', '10,500 - 16,000'],
    ['Legal, Compliance & CA', '8,000 - 10,000'],
    ['Field Onboarding Executive (1x)', '20,000 - 25,000'],
    ['Total Monthly Burn', '43,000 - 63,000'],
]
add_table(s, Inches(0.5), Inches(4.3), Inches(5.5), Inches(3), burn_headers, burn_rows)

# Milestones
add_textbox(s, Inches(6.3), Inches(1.2), Inches(6.5), Inches(0.5),
            'Growth Milestones', font_size=18, bold=True, color=ACCENT_BLUE)

milestones = [
    ('Milestone 1', '50 Paid Sellers', 'Month 4-5', '~Rs 30K MRR', 'Covers basic cloud + amortised setup'),
    ('Milestone 2', '150 Paid Sellers', 'Month 6', '~Rs 1.15L MRR', 'SELF-SUSTAINING — covers full monthly burn'),
    ('Milestone 3', '400 Paid Sellers', 'Month 12', '~Rs 3.75L MRR', 'Solid operational profitability'),
    ('Milestone 4', '2,000 Paid Sellers', 'Year 2', '~Rs 17L MRR', '~Rs 2 Cr ARR — ready for multi-city expansion'),
]

for i, (m, sellers, timeline, mrr, note) in enumerate(milestones):
    y = Inches(1.9) + Inches(i * 1.3)
    color = [ACCENT_RED, ACCENT_GREEN, ACCENT_BLUE, ACCENT_GOLD][i]
    add_shape_rect(s, Inches(6.3), y, Inches(6.5), Inches(1.1), WHITE)
    add_shape_rect(s, Inches(6.3), y, Inches(0.08), Inches(1.1), color)
    add_textbox(s, Inches(6.5), y + Inches(0.05), Inches(2), Inches(0.35),
                f'{m}: {sellers}', font_size=12, bold=True, color=color)
    add_textbox(s, Inches(8.5), y + Inches(0.05), Inches(1.5), Inches(0.35),
                timeline, font_size=11, color=MID_GRAY)
    add_textbox(s, Inches(10), y + Inches(0.05), Inches(2.5), Inches(0.35),
                mrr, font_size=12, bold=True, color=ACCENT_GREEN)
    add_textbox(s, Inches(6.5), y + Inches(0.45), Inches(6), Inches(0.5),
                note, font_size=10, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 9: PRD OVERVIEW
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_accent_bar(s)
add_textbox(s, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            'Product Requirements — 103 Requirements Across 10 Categories', font_size=28, bold=True, color=ACCENT_BLUE)

# KPI boxes
req_kpis = [
    ('103', 'Total Requirements', Inches(0.5)),
    ('13', 'User Workflows', Inches(3.7)),
    ('65', 'Use Cases', Inches(6.9)),
    ('5', 'User Personas', Inches(10.1)),
]
for num, label, left in req_kpis:
    add_kpi_box(s, left, Inches(1.2), Inches(2.8), Inches(1.1), num, label, num_color=ACCENT_BLUE)

# Bar chart for category distribution
cats = ['UI/UX', 'Backend/Data', 'API', 'Infra/Security', 'Commerce', 'Maintenance', 'DR', 'Funnels', 'Perf/Scale', 'Privacy']
vals = [24, 26, 13, 11, 10, 5, 4, 4, 3, 3]
add_chart_bar(s, Inches(0.3), Inches(2.6), Inches(7.5), Inches(4.5), cats, vals, 'Requirements', ACCENT_BLUE)

# Priority breakdown
add_textbox(s, Inches(8), Inches(2.6), Inches(5), Inches(0.4),
            'By Priority Level', font_size=14, bold=True, color=DARK_TEXT)

pri_headers = ['Priority', 'Count', '% of Total']
pri_rows = [
    ['P0 - Critical', '42', '40.8%'],
    ['P1 - High', '41', '39.8%'],
    ['P2 - Medium', '16', '15.5%'],
    ['P3 - Low', '4', '3.9%'],
]
add_table(s, Inches(8), Inches(3.1), Inches(4.8), Inches(2.2), pri_headers, pri_rows)

add_multiline_textbox(s, Inches(8), Inches(5.5), Inches(5), Inches(1.5), [
    'P0 + P1 = 83 requirements (80.6%)',
    'are Critical or High priority',
    '',
    'Categories with most requirements:',
    'Backend/Data: 26 | UI/UX: 24 | API: 13',
    'Infra/Security: 11 | Commerce: 10',
], font_size=11, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 10: USER PERSONAS
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_accent_bar(s)
add_textbox(s, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            'User Personas — Target Audience', font_size=32, bold=True, color=ACCENT_BLUE)

personas = [
    ('Priya Sharma\nHyperlocal Customer', '34, Homemaker &\nBoutique Owner\nBasavanagudi', 'Android & Laptop\nModerate tech literacy', 'Browse + buy across\nmultiple Pete markets\nConsolidated delivery'),
    ('Ramesh Gupta\nTraditional Merchant', '52, 3rd Gen Silk\nShowroom Owner\nChickpet', 'Smartphone\nLow-moderate tech', 'Digital storefront\nWhatsApp enquiries\nBulk order mgmt'),
    ('Vinay Kumar\nDelivery Partner', '28, Independent\nCourier Partner\nNear Chickpet Hub', 'Android Courier App\nModerate tech', 'Multi-store pickup\nMicro-hub consolidation\nSingle-drop delivery'),
    ('Ananya Rao\nPlatform Admin', '31, Operations\nManager\nRemote, Bangalore', 'Laptop + Alerts\nHigh tech literacy', 'Merchant management\nDashboard & analytics\nDispute resolution'),
    ('Deepa Patel\nB2B Wholesale Buyer', '38, Wedding\nPlanner / Reseller\nMalleswaram', 'Smartphone + Laptop\nModerate tech', 'Bulk ordering\nMOQ & tiered pricing\nNegotiation via WhatsApp'),
]

for i, (name, bio, tech, needs) in enumerate(personas):
    x = Inches(0.3) + Inches(i * 2.55)
    box = add_shape_rect(s, x, Inches(1.3), Inches(2.4), Inches(5.5), WHITE)
    add_shape_rect(s, x, Inches(1.3), Inches(2.4), Inches(0.6), ACCENT_BLUE)
    add_textbox(s, x + Inches(0.1), Inches(1.35), Inches(2.2), Inches(0.55),
                name, font_size=10, bold=True, color=WHITE)
    add_textbox(s, x + Inches(0.1), Inches(2.1), Inches(2.2), Inches(0.8),
                bio, font_size=9, color=DARK_TEXT)
    add_textbox(s, x + Inches(0.1), Inches(3.0), Inches(2.2), Inches(0.2),
                'DEVICE & TECH', font_size=7, bold=True, color=MID_GRAY)
    add_textbox(s, x + Inches(0.1), Inches(3.2), Inches(2.2), Inches(0.6),
                tech, font_size=9, color=DARK_TEXT)
    add_textbox(s, x + Inches(0.1), Inches(3.9), Inches(2.2), Inches(0.2),
                'KEY NEEDS', font_size=7, bold=True, color=MID_GRAY)
    add_textbox(s, x + Inches(0.1), Inches(4.1), Inches(2.2), Inches(1.5),
                needs, font_size=9, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 11: SAMPLE REQUIREMENTS
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_accent_bar(s)
add_textbox(s, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            'Key Requirements — Representative Sample', font_size=32, bold=True, color=ACCENT_BLUE)

req_headers = ['Req ID', 'Category', 'Requirement', 'Priority']
req_rows = [
    ['UI-001', 'UI/UX', 'Responsive web app with mobile-first design supporting all 3 modes', 'P0'],
    ['UI-005', 'UI/UX', 'Merchant dashboard for product catalog, order mgmt, analytics', 'P0'],
    ['UI-012', 'UI/UX', 'Multi-store cart with consolidated checkout across Pete markets', 'P0'],
    ['API-001', 'API', 'RESTful API for product catalog CRUD with pagination & filtering', 'P0'],
    ['API-005', 'API', 'WhatsApp Business API integration for Mode B enquiry routing', 'P0'],
    ['BE-001', 'Backend', 'Multi-tenant merchant architecture with isolated store profiles', 'P0'],
    ['BE-005', 'Backend', 'Zone-based dynamic delivery fee calculator with weight surcharges', 'P1'],
    ['BE-012', 'Backend', 'Order lifecycle engine with multi-store pickup routing to micro-hub', 'P0'],
    ['COM-001', 'Commerce', 'Razorpay payment gateway with escrow and split payments', 'P0'],
    ['COM-005', 'Commerce', 'Subscription billing engine with monthly/annual plans', 'P1'],
    ['INFRA-001', 'Infra', 'Supabase PostgreSQL with Row Level Security per merchant', 'P0'],
    ['INFRA-005', 'Infra', 'Redis caching layer for product catalog with TTL invalidation', 'P1'],
    ['SEC-001', 'Security', 'JWT-based authentication with OTP login and role-based access', 'P0'],
    ['DR-001', 'DR', 'Automated daily PostgreSQL backups with 30-day retention', 'P1'],
    ['AI-001', 'AI', 'Auto-generation of product reels from merchant uploaded images', 'P2'],
]
add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5), req_headers, req_rows,
          col_widths=[Inches(1), Inches(1.5), Inches(8), Inches(0.8)])

# ════════════════════════════════════════
# SLIDE 12: MVP SCOPE
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_accent_bar(s)
add_textbox(s, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            'MVP Scope — Phase 1: Core Launch', font_size=32, bold=True, color=ACCENT_BLUE)

add_multiline_textbox(s, Inches(0.5), Inches(1.2), Inches(6), Inches(3), [
    'MVP Includes (P0 - 42 Critical Requirements):',
    '',
    'In-Scope for MVP:',
    '  Web & Mobile (React + React Native)',
    '  Mode A, B, C implementation',
    '  Multi-store cart & consolidated checkout',
    '  Razorpay payment integration',
    '  Merchant onboarding & dashboard',
    '  Product catalog management',
    '  WhatsApp enquiry routing',
    '  Google Maps store integration',
    '  Zone-based hyperlocal delivery',
    '  Admin panel & analytics',
    '  JWT auth + role-based access',
    '  Supabase PostgreSQL + Redis cache',
    '  Basic SEO & Google indexation',
], font_size=11, color=DARK_TEXT)

add_multiline_textbox(s, Inches(6.8), Inches(1.2), Inches(6), Inches(3), [
    'Post-MVP (P1-P3 - 61 requirements):',
    '',
    'Deferred to Post-Launch:',
    '  AI-Generated product reels',
    '  WhatsApp broadcast marketing',
    '  Sponsored CPC listings',
    '  Multi-city expansion framework',
    '  Advanced analytics dashboard',
    '  Customer review system',
    '  Virtual Walk & Live Bazaar (Phase 2)',
    '  Co-Shopping features',
    '  ShipRocket logistics integration',
    '  Data licensing marketplace',
    '  AI Virtual Try-On (jewellery)',
    '  Kill-switch & circuit breaker',
], font_size=11, color=DARK_TEXT)

add_shape_rect(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), LIGHT_GRAY)
add_textbox(s, Inches(0.7), Inches(4.6), Inches(12), Inches(0.4),
            'MVP Success Criteria', font_size=14, bold=True, color=ACCENT_BLUE)
add_multiline_textbox(s, Inches(0.7), Inches(5.1), Inches(11.5), Inches(1.5), [
    '• 50 pilot merchants onboarded across 5 core Pete markets (Chickpet, Balepet, Mamulpet, KR Market, Avenue Road)',
    '• All 3 interaction modes functional with end-to-end order lifecycle',
    '• Multi-store consolidated cart with zone-based delivery working in production',
    '• Monthly Active Merchants (MAM) >= 40 | Monthly Active Buyers (MAB) >= 500 | Avg Order Value >= Rs 800',
    '• Platform Uptime >= 99.5% | API Response Time < 200ms p95',
], font_size=11, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 13: TECH STACK OVERVIEW
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_accent_bar(s)
add_textbox(s, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            'Proposed Technology Stack', font_size=32, bold=True, color=ACCENT_BLUE)

stack_data = [
    ['Frontend (Web)', 'React / Next.js', 'Responsive PWA with SSR, SEO optimization'],
    ['Mobile (Android)', 'React Native', 'Cross-platform shared logic with web'],
    ['Mobile (iOS)', 'React Native', 'Single codebase with Android'],
    ['Backend / API', 'Node.js / Express', 'RESTful API gateway with rate limiting'],
    ['Database', 'Supabase (PostgreSQL)', 'Multi-tenant with Row Level Security'],
    ['Cache', 'Redis', 'Product catalog cache with TTL invalidation'],
    ['Auth', 'Supabase Auth + JWT', 'OTP login, role-based access, SSO ready'],
    ['Payments', 'Razorpay', 'Escrow, split payments, subscription billing'],
    ['Hosting', 'Vercel + Railway', 'Serverless frontend + containerized backend'],
    ['Messaging', 'WhatsApp Business API', 'Deep-link enquiries, order notifications'],
    ['Storage', 'Supabase Storage', 'Product images, store photos, reels'],
    ['Search', 'Elasticsearch / MeiliSearch', 'Full-text product & store search'],
    ['Analytics', 'PostHog + Metabase', 'Product analytics + business intelligence'],
]

stack_headers = ['Layer', 'Technology', 'Purpose']
add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5), stack_headers, stack_data,
          col_widths=[Inches(2.5), Inches(3), Inches(6.8)])

# ════════════════════════════════════════
# SLIDE 14: ROADMAP
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_accent_bar(s)
add_textbox(s, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
            'Implementation Roadmap & Next Steps', font_size=32, bold=True, color=ACCENT_BLUE)

phases = [
    ('PHASE 1\nFoundation', 'Month 1-2', ACCENT_RED, Inches(0.5), [
        'Tech stack finalisation',
        'Architecture blueprint',
        'POC & prototype build',
        'DevOps pipeline setup',
        'Sprint 0: Foundation',
    ]),
    ('PHASE 2\nCore Build', 'Month 3-4', ORANGE, Inches(3.5), [
        'MVP frontend (web + mobile)',
        'API gateway & endpoints',
        'Database schema & migration',
        'Integration & QA',
        'Staging deployment',
    ]),
    ('PHASE 3\nLaunch', 'Month 5-6', ACCENT_GREEN, Inches(6.5), [
        '50 pilot merchant onboarding',
        'Production deployment',
        'Marketing & SEO launch',
        'Customer support setup',
        'Post-launch monitoring',
    ]),
    ('PHASE 4\nScale', 'Month 7-12', ACCENT_BLUE, Inches(9.5), [
        '150+ merchant acquisition',
        'AI reels & smart features',
        'Multi-city expansion prep',
        'B2B wholesale marketplace',
        'Target: Rs 3.75L MRR',
    ]),
]

for title, timeline, color, left, items in phases:
    box = add_shape_rect(s, left, Inches(1.3), Inches(2.8), Inches(5.5), WHITE)
    add_shape_rect(s, left, Inches(1.3), Inches(2.8), Inches(1.0), color)
    add_textbox(s, left + Inches(0.15), Inches(1.35), Inches(2.5), Inches(0.6),
                title, font_size=12, bold=True, color=WHITE)
    add_textbox(s, left + Inches(0.15), Inches(1.85), Inches(2.5), Inches(0.3),
                timeline, font_size=10, color=RGBColor(0xDD, 0xDD, 0xDD))
    add_multiline_textbox(s, left + Inches(0.15), Inches(2.6), Inches(2.5), Inches(3.5), items, font_size=10, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 15: CLOSING
# ════════════════════════════════════════
s = new_slide()
add_bg(s, BG_DARK)
add_textbox(s, Inches(0.8), Inches(1.5), Inches(11), Inches(1),
            'PeteMart', font_size=54, bold=True, color=ACCENT_GOLD)
add_textbox(s, Inches(0.8), Inches(2.7), Inches(11), Inches(0.8),
            'Transforming Old Bangalore\'s Historic Trade into India\'s Most Unique E-Commerce Experience',
            font_size=18, color=WHITE)
add_shape_rect(s, Inches(0.8), Inches(3.7), Inches(3), Inches(0.04), ACCENT_GOLD)
add_multiline_textbox(s, Inches(0.8), Inches(4.0), Inches(11), Inches(2), [
    '103 Requirements Ready | 406 Merchants Profiled | 21 Markets Mapped',
    'Setup Cost: Rs 86K - 1.5L | Monthly Burn: Rs 43K - 63K | Break-Even: Month 6',
    '',
    'Next Step: Architecture Design & POC — Awaiting Management Approval',
], font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

add_textbox(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
            'Thank You | Questions & Discussion', font_size=16, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

# ── Save ──
output_path = os.path.join(os.getcwd(), '01_front_office', '01_ideation_agent', '01_ideation_agent_FULL_PRESENTATION.pptx')
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
