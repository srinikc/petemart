from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_SLIDE = RGBColor(0xF5, 0xF0, 0xEB)
ACCENT_GOLD = RGBColor(0xD4, 0xA0, 0x3C)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
ACCENT_BLUE = RGBColor(0x15, 0x63, 0x8C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
LIGHT_GRAY = RGBColor(0xF0, 0xEE, 0xE9)
MID_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
TEAL = RGBColor(0x00, 0x7B, 0x7F)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper Functions ──
def add_bg(slide, color=BG_SLIDE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_rounded_rect(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_accent_bar(slide, left=Inches(0), top=Inches(0), width=Inches(0.15), height=Inches(7.5), color=ACCENT_GOLD):
    return add_shape_rect(slide, left, top, width, height, color)

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tb

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=12, color=DARK_TEXT, bold_first=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        if bold_first and i == 0:
            p.font.bold = True
            p.font.size = Pt(font_size + 2)
        p.space_after = Pt(3)
    return tb

def add_kpi_box(slide, left, top, width, height, number, label, bg_color=WHITE, num_color=ACCENT_GOLD):
    add_shape_rect(slide, left, top, width, height, bg_color)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.6),
                str(number), font_size=28, bold=True, color=num_color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.65), width - Inches(0.3), Inches(0.4),
                label, font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE
            p.font.name = 'Calibri'
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = DARK_TEXT
                p.font.name = 'Calibri'
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 0 else LIGHT_GRAY
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    return tbl

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_section_header(slide, text, subtitle=None):
    add_accent_bar(slide)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
                text, font_size=32, bold=True, color=ACCENT_BLUE)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.5),
                    subtitle, font_size=13, color=MID_GRAY)

def add_architecture_box(slide, left, top, width, height, title, items, color=ACCENT_BLUE):
    add_shape_rect(slide, left, top, width, height, WHITE)
    add_shape_rect(slide, left, top, width, Inches(0.45), color)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.05), width - Inches(0.3), Inches(0.35),
                title, font_size=11, bold=True, color=WHITE)
    add_multiline_textbox(slide, left + Inches(0.15), top + Inches(0.55), width - Inches(0.3), height - Inches(0.6),
                          items, font_size=9, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════
s = new_slide()
add_bg(s, BG_DARK)
add_textbox(s, Inches(0.8), Inches(0.6), Inches(11), Inches(1),
            'PeteMart', font_size=54, bold=True, color=ACCENT_GOLD)
add_textbox(s, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
            'Enterprise Architecture Blueprint — Dual Architecture Design Review', font_size=22, color=WHITE)
add_shape_rect(s, Inches(0.8), Inches(2.5), Inches(3), Inches(0.04), ACCENT_GOLD)
add_multiline_textbox(s, Inches(0.8), Inches(2.8), Inches(11), Inches(3), [
    'Architecture Review Package',
    'Agent 03 — Senior Enterprise Solution Architect',
    '',
    'Presented for: Higher Management / CTO Review',
    'Date: May 2026',
    'Coverage: Full Product Architecture + Zero-Cost POC Blueprint',
    'Status: 103 Requirements Mapped | All Guardrails Passed ✅',
], font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

# ════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Executive Summary — Dual Architecture Strategy',
                   'Two blueprints: Full Product (5,000+ merchants) + POC (8 merchants, ₹0/month)')

kpis = [
    ('103', 'Requirements Mapped', Inches(0.5), ACCENT_BLUE),
    ('48', 'POC Requirements', Inches(3.7), ACCENT_GREEN),
    ('35', 'Database Tables', Inches(6.9), PURPLE),
    ('₹0/mo', 'POC Cloud Cost', Inches(10.1), ACCENT_GOLD),
]
for num, label, left, color in kpis:
    add_kpi_box(s, left, Inches(1.3), Inches(2.8), Inches(1.1), num, label, num_color=color)

add_textbox(s, Inches(0.5), Inches(2.7), Inches(6), Inches(0.5),
            'Architecture Principles', font_size=18, bold=True, color=ACCENT_BLUE)
add_multiline_textbox(s, Inches(0.5), Inches(3.3), Inches(6), Inches(3.5), [
    'API-First: All frontends consume same RESTful API (/api/v1/*)',
    '  Web (Next.js), Mobile (Expo RN), Third-party integrations',
    '',
    'Serverless-First: Supabase (Postgres/Auth/Storage/Realtime)',
    '  eliminates dedicated backend servers in POC; scales to production',
    '',
    'Zero-Cost POC: Only free tiers — Supabase Free, Vercel Hobby,',
    '  Railway $5 credit, GitHub Pages, Expo Go',
    '',
    'AI-Enhanced UI: Google Stitch is PRIMARY UI design tool —',
    '  generates React code from NL prompts (@google/stitch-sdk)',
    '',
    'Multi-Tenancy: Each merchant gets isolated microsite',
    '  with white-label theming and Row Level Security',
], font_size=11, color=DARK_TEXT)

add_textbox(s, Inches(6.8), Inches(2.7), Inches(6), Inches(0.5),
            'Key Decisions', font_size=18, bold=True, color=ACCENT_BLUE)
add_multiline_textbox(s, Inches(6.8), Inches(3.3), Inches(6), Inches(3.5), [
    'Web-first, mobile later (PWA + Expo Go for POC)',
    'Supabase PostgreSQL > Firestore/MongoDB (RLS + maturity)',
    'Google Stitch > Figma/Sketch (AI-native generation)',
    'Razorpay > Stripe/Cashfree (India-dominant + subaccount)',
    'PgFTS > Algolia/Elasticsearch (free, sufficient for POC)',
    'WhatsApp Business API > Twilio (free tier + India ubiquity)',
    'Next.js SSR > SPA (SEO for merchant microsites)',
    'No i18n for POC (English + hardcoded Kannada)',
    'Single admin for POC (full RBAC in production)',
    'Manual merchant onboarding for pilot (automation later)',
], font_size=11, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 3: SYSTEM CONTEXT (C4 LEVEL 1)
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'System Context Diagram — C4 Level 1',
                   'PeteMart platform interacting with 5 user types and 7 external systems')

# External systems as boxes
systems = [
    ('Users (5)', ACCENT_BLUE, Inches(0.5), [
        'Customer / B2B Buyer',
        'Pete Merchant',
        'Delivery Partner',
        'Platform Admin',
        'UI Designer (Agent 07a)',
    ]),
    ('PeteMart Platform', ACCENT_GREEN, Inches(4.6), [
        'Web App (Next.js 15 + shadcn/ui)',
        'Mobile App (Expo React Native)',
        'Courier App (Expo React Native)',
        'Supabase (Postgres + Auth + Storage + Realtime)',
        'API Gateway (Next.js Middleware)',
    ]),
    ('External Systems (7)', ORANGE, Inches(8.7), [
        'Razorpay — Payment Gateway',
        'WhatsApp Business API — Enquiry/Notifications',
        'Google Maps — Geocoding/Directions',
        'Google Gemini — AI Inference (Try-On, Moderation)',
        'Google Stitch — AI UI Generation',
        'ShipRocket — National Shipping Labels',
        'Bullion Rate API — Live Gold/Silver Rates',
    ]),
]
for title, color, left, items in systems:
    add_architecture_box(s, left, Inches(1.4), Inches(3.8), Inches(5), title, items, color)

add_textbox(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
            'All diagrams rendered in Mermaid.js — viewable on GitHub, Notion, or mermaid.live',
            font_size=10, color=MID_GRAY)

# ════════════════════════════════════════
# SLIDE 4: CONTAINER DIAGRAM (C4 LEVEL 2)
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Container Diagram — C4 Level 2',
                   'All containers with technology stack and interactions')

layers = [
    ('Presentation Layer', ACCENT_BLUE, Inches(0.3), [
        'Next.js 15 + shadcn/ui + Tailwind',
        '  → Vercel Hobby (POC) / Pro (Production)',
        '  → SSR, PWA, Service Worker',
        'Expo React Native + NativeWind',
        '  → Expo Go (POC) / EAS Build (Prod)',
        'Google Stitch (@google/stitch-sdk)',
        '  → 350 std + 200 pro generations/mo FREE',
    ]),
    ('API Layer', PURPLE, Inches(4.6), [
        'Next.js API Routes (/api/v1/*)',
        '  → Rate limiting (Token Bucket)',
        '  → Zod validation + JWT auth middleware',
        'Supabase Edge Functions (Deno)',
        '  → 500K invocations/mo FREE',
        '  → Payment webhooks, delivery logic',
    ]),
    ('Data Layer', ACCENT_GREEN, Inches(8.9), [
        'Supabase PostgreSQL (35 tables)',
        '  → 500MB FREE / 8GB Pro',
        '  → Row Level Security per merchant',
        'Redis Cache (Upstash Free)',
        '  → 100MB, product catalog TTL',
        'Supabase Storage (S3-compatible)',
        '  → 1GB FREE / 100GB + Pro',
        'Supabase Realtime (WebSocket)',
        '  → 200 concurrent FREE',
    ]),
]
for title, color, left, items in layers:
    add_architecture_box(s, left, Inches(1.3), Inches(4.1), Inches(5.2), title, items, color)

add_textbox(s, Inches(0.3), Inches(6.7), Inches(12), Inches(0.5),
            '35 DB tables → 15 in POC | All frontends consume same REST API | Caching: Redis + Supabase query cache',
            font_size=10, color=MID_GRAY)

# ════════════════════════════════════════
# SLIDE 5: TECH STACK & RATIONALE
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Technology Stack — Decision Matrix',
                   'Every layer evaluated against alternatives with clear rationale')

stack_headers = ['Layer', 'Chosen', 'Alternatives', 'Rationale']
stack_rows = [
    ['UI Design', 'Google Stitch', 'Figma, Sketch', 'AI-native NL→UI generation; FREE; React code export'],
    ['Web Frontend', 'Next.js 15 / shadcn/ui', 'Nuxt, SvelteKit, SPA', 'SSR for SEO (microsites); Vercel-native; accessible'],
    ['Mobile', 'Expo React Native', 'Flutter, Swift/Kotlin', 'Expo Go free testing; shared web logic; NativeWind'],
    ['Database', 'Supabase PostgreSQL', 'Firestore, MongoDB', 'Maturity + RLS + Auth + Storage + Realtime ($0)'],
    ['Backend', 'Edge Functions + API Routes', 'Express, FastAPI', 'Serverless; co-located with DB; Deno runtime'],
    ['Auth', 'Supabase Auth', 'Auth0, Clerk', 'Built-in; OTP/phone; JWT; RLS integration'],
    ['Payments', 'Razorpay', 'Stripe, Cashfree', 'India-dominant; test mode free; subaccount Route API'],
    ['Cache', 'Upstash Redis', 'ElastiCache, self-hosted', '$9/mo production; free tier for POC'],
    ['Search', 'PgFTS → Meilisearch', 'Algolia, Elasticsearch', 'PgFTS free; Meilisearch $29/mo for production'],
    ['AI', 'Gemini (Free) → Pro', 'OpenAI, Claude', '60 req/min free; multilingual; Stitch-native'],
    ['CI/CD', 'GitHub Actions', 'GitLab CI, CircleCI', 'Free for public repos; tight GitHub integration'],
    ['Monitoring', 'Sentry + Vercel Analytics', 'Datadog, NewRelic', 'Sentry free (5K errors); built-in Vercel RUM'],
    ['Messaging', 'WhatsApp Business API', 'Twilio, MSG91', '1K convos/mo free; India ubiquity; deep links'],
    ['Hosting', 'Vercel + Supabase', 'AWS, GCP, Azure', 'Vercel Hobby free; Supabase free; serverless-first'],
]
add_table(s, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5.8), stack_headers, stack_rows,
          col_widths=[Inches(1.5), Inches(2.5), Inches(2.5), Inches(6.2)])

# ════════════════════════════════════════
# SLIDE 6: GOOGLE STITCH INTEGRATION
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Google Stitch — AI-Powered UI Generation Pipeline',
                   'PRIMARY UI design tool for PeteMart. Generates React code from natural language prompts.')

add_kpi_box(s, Inches(0.5), Inches(1.3), Inches(2), Inches(1.1), '550/mo', 'FREE Generations', num_color=ACCENT_GREEN)
add_kpi_box(s, Inches(2.8), Inches(1.3), Inches(2), Inches(1.1), '30+', 'Screens Generated', num_color=ACCENT_BLUE)
add_kpi_box(s, Inches(5.1), Inches(1.3), Inches(2), Inches(1.1), 'MCP SDK', 'Programmatic API', num_color=PURPLE)
add_kpi_box(s, Inches(7.4), Inches(1.3), Inches(2.5), Inches(1.1), 'DESIGN.md', 'Token Handoff', num_color=ACCENT_GOLD)

add_multiline_textbox(s, Inches(0.5), Inches(2.7), Inches(6), Inches(4), [
    'Why Google Stitch is PRIMARY UI Tool:',
    '',
    '• AI-Native Design: Unlike Figma (manual) or Sketch (static),',
    '  Stitch generates UI from natural language prompts',
    '',
    '• MCP SDK Integration (@google/stitch-sdk):',
    '    stitch.generate() — React components from NL prompts',
    '    stitch.design() — Full design system tokens',
    '    stitch.export() — React code + HTML + DESIGN.md',
    '    stitch.prototype() — Interactive multi-screen prototypes',
    '',
    '• Zero-Cost: 350 standard + 200 pro gens/month FREE',
    '  (~40 gens needed for entire POC)',
    '',
    '• DESIGN.md Workflow: Exports color palette, typography,',
    '  spacing grid, component variants, breakpoints',
], font_size=11, color=DARK_TEXT)

add_multiline_textbox(s, Inches(6.8), Inches(2.7), Inches(6), Inches(4), [
    'Stitch Integration Pipeline:',
    '',
    'Agent 07a (UI) writes NL prompts →',
    'Stitch SDK generates React .tsx →',
    'DESIGN.md tokens exported →',
    'Tailwind config auto-updated →',
    'Components integrated into Next.js app →',
    'Screenshots generated for human review →',
    '',
    'Stitch-generated screens for POC:',
    '• Landing Page w/ Pete Tapestry (2 gens)',
    '• Product Catalog & Search (2 gens)',
    '• Product Cards (Mode A/B/C badges) (2 gens)',
    '• Multi-Store Cart & Checkout (2 gens)',
    '• Merchant Store Microsite (2 gens)',
    '• Merchant Dashboard (2 gens)',
    '• Admin Console (2 gens)',
    '• Order Tracking Dashboard (1 gen)',
    '• Review Interface (1 gen)',
    '',
    'Voice Canvas also supported for accessibility.',
], font_size=11, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 7: DATA FLOW — ORDER LIFECYCLE
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Data Flow — End-to-End Order Lifecycle',
                   'Mode A (Direct Purchase) full flow from Browse to Settlement')

order_flow = [
    ('1. Browse & Cart', ACCENT_BLUE, [
        'Customer browses products across Pete markets',
        'Adds items from Merchant A + Merchant B to cart',
        'Cart service validates stock & pricing',
        'Multi-store consolidation flag set',
    ]),
    ('2. Checkout & Fee', ACCENT_GREEN, [
        'Delivery fee calculator (zone + weight + consolidation)',
        'Example: 3 merchants, Zone 2 → ₹130 total',
        'Razorpay order created via Route API',
    ]),
    ('3. Payment', ACCENT_GOLD, [
        'Customer pays via UPI / Card / Wallet',
        'Razorpay webhook: payment.captured',
        'Signature + amount verified server-side',
    ]),
    ('4. Fulfillment', ORANGE, [
        'Order split by merchant (consolidation_id)',
        'Courier assigned (zone-based)',
        'Notifications sent to all parties',
    ]),
    ('5. Pickup & Hub', PURPLE, [
        'Courier picks up from Merchant A → Merchant B',
        'Items consolidated at Chickpet micro-hub',
        'GPS location broadcast (30s intervals)',
    ]),
    ('6. Delivery & Settle', ACCENT_RED, [
        'Last-mile delivery to customer',
        'Order status → DELIVERED',
        'Settlement timer T+3 → Razorpay Route API payout',
        'Customer prompted to rate',
    ]),
]
for i, (title, color, items) in enumerate(order_flow):
    x = Inches(0.3) + Inches(i * 2.15)
    y = Inches(1.3)
    w = Inches(2.0)
    h = Inches(5.5)
    add_shape_rect(s, x, y, w, h, WHITE)
    add_shape_rect(s, x, y, w, Inches(0.45), color)
    add_textbox(s, x + Inches(0.08), y + Inches(0.05), w - Inches(0.16), Inches(0.35),
                title, font_size=10, bold=True, color=WHITE)
    add_multiline_textbox(s, x + Inches(0.08), y + Inches(0.55), w - Inches(0.16), h - Inches(0.6),
                          items, font_size=8, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 8: API STRATEGY
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'API Strategy & Rate Limiting',
                   'API-First: All frontends (Web, Mobile, Courier) consume same REST API with token bucket rate limiting')

api_headers = ['Endpoint', 'Method', 'Rate Limit', 'Auth', 'Description']
api_rows = [
    ['/api/v1/products', 'GET/POST/PUT', '100/min per IP', 'Public + JWT', 'Product catalog CRUD with pagination'],
    ['/api/v1/auth/*', 'POST', '10/min per phone', 'None', 'Phone OTP login + JWT generation'],
    ['/api/v1/orders', 'GET/POST/PUT', '60/min per user', 'JWT', 'Order lifecycle management'],
    ['/api/v1/cart', 'GET/POST/PUT/DELETE', '60/min per user', 'JWT', 'Multi-store cart operations'],
    ['/api/v1/checkout', 'POST', '30/min per user', 'JWT', 'Checkout + delivery fee calculation'],
    ['/api/v1/merchants', 'GET/POST/PUT', '60/min per IP', 'JWT (Admin)', 'Merchant profile management'],
    ['/api/v1/whatsapp/deeplink', 'GET', '30/min per IP', 'Public', 'WhatsApp deep link generation'],
    ['/api/v1/directions', 'GET', '100/min per IP', 'Public', 'Google Maps directions + geocoding'],
    ['/api/v1/tracking', 'GET/PUT', '300/min per courier', 'JWT', 'GPS tracking (30s intervals)'],
    ['/api/v1/admin/*', 'GET/POST/PUT', '60/min per admin', 'JWT (Admin)', 'Admin console operations'],
    ['/api/v1/bullion', 'GET', '60/min per IP', 'Public+cache', 'Live gold/silver rates (cached 5min)'],
    ['/api/v1/notifications', 'POST', '30/min per user', 'JWT', 'Send SMS/email/push notifications'],
]
add_table(s, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5), api_headers, api_rows,
          col_widths=[Inches(2.8), Inches(1.2), Inches(2), Inches(1.5), Inches(5.2)])

add_textbox(s, Inches(0.3), Inches(6.5), Inches(12), Inches(0.5),
            'Rate Limiting Algorithm: Token Bucket | Implementation: Next.js Middleware + Zod Validation | POC: Basic rate limiting via middleware',
            font_size=10, color=MID_GRAY)

# ════════════════════════════════════════
# SLIDE 9: SECURITY ARCHITECTURE
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Security Architecture — 10-Layer Defense',
                   'Comprehensive security framework across all application layers')

sec_layers = [
    ('Layer 1: Transport', ACCENT_RED, [
        'HTTPS/TLS 1.3 (Vercel + Supabase)',
        'HSTS headers enforced',
        'Certificate auto-renewal (Let\'s Encrypt)',
    ]),
    ('Layer 2: Auth', ACCENT_GREEN, [
        'Supabase Auth (Phone OTP + JWT)',
        'JWT expiry: 1hr access, 30d refresh',
        'Role-based: Customer / Merchant / Admin / Courier',
    ]),
    ('Layer 3: API', ACCENT_BLUE, [
        'Rate limiting (Token Bucket)',
        'Request validation (Zod schemas)',
        'CORS whitelist (petemart domains)',
        'SQL injection prevention (parameterized queries)',
    ]),
    ('Layer 4: Database', PURPLE, [
        'Row Level Security per merchant',
        'RLS policies: SELECT/INSERT/UPDATE/DELETE',
        'Encrypted at rest (Supabase managed)',
    ]),
    ('Layer 5: Storage', ORANGE, [
        'RLS on storage buckets',
        'Image transformation URLs signed',
        'Public bucket for store images only',
    ]),
    ('Layer 6: Payments', ACCENT_GOLD, [
        'Razorpay webhook signature verification',
        'Amount validation server-side',
        'Test mode for POC (no real transactions)',
    ]),
    ('Layer 7: Secrets', ACCENT_RED, [
        'All secrets in Vercel env vars (encrypted)',
        'No .env committed to git',
        'Secrets rotation policy documented',
    ]),
    ('Layer 8: Compliance', ACCENT_GREEN, [
        'GST invoice generation (BE-015)',
        'TCS compliance for e-commerce',
        'Basic privacy policy (GDPR-readiness)',
    ]),
    ('Layer 9: Monitoring', ACCENT_BLUE, [
        'Sentry error tracking (5K errors free)',
        'Vercel Analytics (RUM + Web Vitals)',
        'Supabase Dashboard (DB monitoring)',
    ]),
    ('Layer 10: Testing', PURPLE, [
        'Dependabot auto-dependency scanning',
        'Snyk vulnerability scanning (Team $25/mo)',
        'OWASP ZAP security scan in CI/CD',
    ]),
]
for i, (title, color, items) in enumerate(sec_layers):
    x = Inches(0.3) + Inches(i * 1.27)
    add_architecture_box(s, x, Inches(1.3), Inches(1.2), Inches(5.5), title, items, color)

# ════════════════════════════════════════
# SLIDE 10: INFRASTRUCTURE COSTING
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Infrastructure Costing — POC vs Full Production',
                   'POC: ₹0/month cloud cost | Production: ₹32,465/month (5,000 merchants)')

# POC Cost
add_shape_rect(s, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.5), WHITE)
add_shape_rect(s, Inches(0.3), Inches(1.3), Inches(6.2), Inches(0.5), ACCENT_GREEN)
add_textbox(s, Inches(0.5), Inches(1.35), Inches(5.8), Inches(0.45),
            'POC Phase — ₹0 / Month Cloud Cost', font_size=14, bold=True, color=WHITE)

poc_cost_headers = ['Service', 'Plan', 'Cost']
poc_cost_rows = [
    ['Vercel Hobby', 'Free', '₹0'],
    ['Supabase Free', 'Free', '₹0'],
    ['Railway', '$5 credit', '₹0 (~₹415 one-time)'],
    ['Expo Go', 'Free', '₹0'],
    ['Google Stitch', 'Free (550 gens)', '₹0'],
    ['Google Gemini', 'Free (60 req/min)', '₹0'],
    ['Razorpay Test', 'Test Mode', '₹0'],
    ['WhatsApp Business', 'Free (1K convos)', '₹0'],
    ['GitHub Actions', 'Free (2K min)', '₹0'],
    ['Sentry Free', 'Free (5K errors)', '₹0'],
    ['Total Cloud Cost', '', '₹0 / MONTH'],
]
add_table(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.5), poc_cost_headers, poc_cost_rows,
          col_widths=[Inches(2.5), Inches(2), Inches(1.3)])

# Production Cost
add_shape_rect(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.5), WHITE)
add_shape_rect(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(0.5), ACCENT_BLUE)
add_textbox(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.45),
            'Full Production — ₹32,465 / Month (5,000 merchants)', font_size=14, bold=True, color=WHITE)

prod_cost_headers = ['Service', 'Tier', 'Monthly (₹)']
prod_cost_rows = [
    ['Vercel Pro (×2)', '$20/mo ×2', '₹3,370'],
    ['Supabase Pro', '$25/mo', '₹2,100'],
    ['Storage Add-on', '$10/100GB', '₹840'],
    ['Upstash Redis Pro', '$9/mo', '₹760'],
    ['Meilisearch Cloud', '$29/mo', '₹2,440'],
    ['Google Maps API', 'Pay-as-you-go', '₹3,000'],
    ['WhatsApp Business', 'Convo-based', '₹3,000'],
    ['SendGrid Essentials', '50K emails', '₹1,200'],
    ['Gemini Pro', 'Usage-based', '₹5,000'],
    ['Sentry Team', '$26/mo', '₹2,200'],
    ['Domain, Monitoring, Security', 'Various', '₹8,555'],
    ['Total Cloud Cost', '', '₹32,465 / MONTH'],
]
add_table(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.5), prod_cost_headers, prod_cost_rows,
          col_widths=[Inches(2.5), Inches(1.8), Inches(1.5)])

# ════════════════════════════════════════
# SLIDE 11: POC SCOPE & PILOT MERCHANTS
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'POC Scope — 8-Merchant Pilot (₹0/month)',
                   '48 of 103 requirements | 9 merchant profiles (8 unique + 1 branch) in Chickpet & Balepet')

add_kpi_box(s, Inches(0.5), Inches(1.3), Inches(1.8), Inches(1.1), '48', 'Requirements Built', num_color=ACCENT_GREEN)
add_kpi_box(s, Inches(2.5), Inches(1.3), Inches(1.8), Inches(1.1), '9', 'Pilot Merchants', num_color=ACCENT_BLUE)
add_kpi_box(s, Inches(4.5), Inches(1.3), Inches(1.8), Inches(1.1), '15', 'DB Tables (POC)', num_color=PURPLE)
add_kpi_box(s, Inches(6.5), Inches(1.3), Inches(2.5), Inches(1.1), '2 Markets', 'Chickpet + Balepet', num_color=ACCENT_GOLD)

m_headers = ['Merchant', 'Market', 'Category', 'Readiness', 'POC Modes']
m_rows = [
    ['Tarun Enterprises', 'Chickpet', 'Textiles Wholesale', 'Low', 'B, C'],
    ['Sri Vari Traders', 'Balepet', 'Outdoor Equipment', 'Low', 'A, B'],
    ['Samskruti Silks (Store 1)', 'Chickpet', 'Silk Sarees', 'Medium', 'A, B, C'],
    ['Samskruti Silks (Branch)', 'Chickpet', 'Silk Sarees', 'Medium', 'A, B, C'],
    ['flowers2u', 'Balepet', 'Florist', 'Low', 'A, B'],
    ['The Pastry Cafe', 'Balepet', 'Bakery & Cafe', 'Medium', 'A, C'],
    ['Sri Vinayaka Textorium', 'Balepet', 'Textiles & Fabrics', 'Low', 'B, C'],
    ['Sanjana Apparels', 'Balepet', 'Apparel & Clothing', 'Medium', 'A, B, C'],
    ['Madhumathi Men\'s Ethnic', 'Balepet', 'Ethnic Wear', 'Low', 'A, B, C'],
]
add_table(s, Inches(0.5), Inches(2.7), Inches(7), Inches(3.5), m_headers, m_rows,
          col_widths=[Inches(2.5), Inches(1.2), Inches(1.5), Inches(0.9), Inches(0.9)])

add_multiline_textbox(s, Inches(7.8), Inches(2.7), Inches(5), Inches(3.5), [
    'Requirements Breakdown (POC):',
    '',
    'UI/UX: 12 of 24 (P0 screens)',
    'API: 8 of 13 (core endpoints)',
    'Backend/Data: 12 of 26 (core schema)',
    'Commerce: 4 of 10 (basic checkout)',
    'Infra/Security: 6 of 11 (free tier)',
    'Performance: 2 of 3 (baseline)',
    'Maintenance: 2 of 5 (basic)',
    'DR: 1 of 4 (Supabase backups)',
    'Funnels: 1 of 4 (PWA)',
    '',
    'POC Success Criteria:',
    '• 9/9 merchants with live store',
    '• ≥20 products per merchant',
    '• ≥10 test orders (Mode A)',
    '• ≥5 WhatsApp clicks (Mode B)',
    '• ≥5 Directions clicks (Mode C)',
    '• LCP < 2.5s | API < 300ms p95',
], font_size=10, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 12: DEPLOYMENT TOPOLOGY
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Deployment Topology — Zero-Cost Free Tier',
                   '100% free deployment: Vercel Hobby + Supabase Free + Railway $5 credit + GitHub Pages')

deploy_boxes = [
    ('Vercel Hobby (FREE)', ACCENT_GREEN, Inches(0.3), [
        'Next.js 15 App (SSR + PWA)',
        'API Routes (/api/v1/*)',
        'Vercel Edge Network (Global CDN)',
        '100GB BW, 100 builds/day',
        'Vercel Analytics (Web Vitals)',
        'Service Worker (Offline Cache)',
        'Domain: petemart-poc.vercel.app',
    ]),
    ('Supabase Free (FREE)', ACCENT_BLUE, Inches(4.6), [
        'PostgreSQL 15 (500MB DB)',
        'Auth (Phone OTP + JWT, 50K users)',
        'Storage (1GB, Image CDN)',
        'Realtime (200 concurrent WebSocket)',
        'Edge Functions (Deno, 500K invocations)',
        'Daily automated backups (included)',
        'Domain: project.supabase.co',
    ]),
    ('Railway + GitHub (FREE)', PURPLE, Inches(8.9), [
        'Railway: $5 credit for cron jobs',
        '  Bullion rate fetcher, session cleanup',
        'GitHub Pages: Architecture docs',
        '  petemart.github.io (static hosting)',
        'GitHub Actions: CI/CD (2K min/mo)',
        'Dependabot: Auto-dependency updates',
        'Expo Go: Mobile app local testing',
    ]),
]
for title, color, left, items in deploy_boxes:
    add_architecture_box(s, left, Inches(1.3), Inches(4.1), Inches(4.5), title, items, color)

add_textbox(s, Inches(0.3), Inches(6.0), Inches(12.5), Inches(1.2), '',
            font_size=10, color=MID_GRAY)
add_shape_rect(s, Inches(0.3), Inches(6.0), Inches(12.5), Inches(1.2), LIGHT_GRAY)
add_multiline_textbox(s, Inches(0.5), Inches(6.1), Inches(12), Inches(1), [
    'External Services: Razorpay Test Mode | WhatsApp Business API (1K convos/mo free) | Google Maps ($200/mo credit) | Google Gemini (60 req/min free) | Sentry (5K errors/mo free)',
    'No custom domains needed — use *.vercel.app, *.supabase.co subdomains. POC cloud cost = ₹0/month. Production: Vercel Pro + Supabase Pro → ₹32,465/mo for 5K merchants.',
], font_size=9, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 13: CACHING & MESSAGE QUEUE
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Caching Strategy & Message Queue Architecture',
                   'Multi-layer caching + event-driven webhook architecture for high concurrency')

cache_areas = [
    ('Caching Layers', ACCENT_GREEN, Inches(0.3), [
        'Layer 1: Browser Cache (Service Worker)',
        '  → Static assets, product images',
        '  → PWA offline support',
        '',
        'Layer 2: Redis (Upstash Free/Pro)',
        '  → Product catalog (TTL: 5 min)',
        '  → Session data (TTL: 30 min)',
        '  → Merchant config (TTL: 15 min)',
        '  → Bullion rates (TTL: 5 min, cached)',
        '',
        'Layer 3: Supabase Query Cache',
        '  → Frequent queries auto-cached',
        '  → Prepared statements for perf',
        '',
        'Layer 4: CDN (Vercel Edge)',
        '  → Static pages, images, API responses',
        '  → Global edge network (100+ locations)',
    ]),
    ('Message Queue / Event Architecture', ACCENT_BLUE, Inches(6.8), [
        'Event-driven via Supabase Realtime:',
        '',
        '  Order Placed →',
        '    • Razorpay payment webhook listener',
        '    • Notification service (SMS/Email/Push)',
        '    • Courier assignment engine',
        '    • Analytics event log',
        '',
        '  Payment Captured →',
        '    • Order status update (CONFIRMED)',
        '    • Inventory deduction',
        '    • Merchant notification',
        '    • Settlement timer start (T+3)',
        '',
        '  Delivery Status Change →',
        '    • GPS broadcast to courier app',
        '    • Customer tracking dashboard update',
        '    • Analytics pipeline trigger',
        '',
        '  Webhook Flow:',
        '    Razorpay → Edge Function → DB Write',
        '    → Realtime broadcast → UI update',
    ]),
]
for title, color, left, items in cache_areas:
    add_architecture_box(s, left, Inches(1.3), Inches(6.2), Inches(5.5), title, items, color)

# ════════════════════════════════════════
# SLIDE 14: MULTI-LAYER TESTING
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Multi-Layer Testing Architecture',
                   '5-layer automated testing with CI/CD gates and coverage thresholds')

test_layers = [
    ('Layer 1\nUnit Tests', ACCENT_GREEN, Inches(0.3), [
        'Framework: Vitest + React Testing Library',
        'Coverage: 80%+ (all components)',
        'Scope: Components, hooks, utils, API routes',
        'Run: On every `git push` (pre-commit hook)',
        'Gate: <80% = BLOCKED',
    ]),
    ('Layer 2\nIntegration Tests', ACCENT_BLUE, Inches(2.8), [
        'Framework: Supertest + MSW (Mock Service Worker)',
        'Coverage: API endpoint contracts',
        'Scope: Request/response validation, auth flows',
        'Run: On PR to develop branch',
        'Gate: Any failure = BLOCKED',
    ]),
    ('Layer 3\nE2E Tests', PURPLE, Inches(5.3), [
        'Framework: Playwright (Web) + Detox (Mobile)',
        'Coverage: Critical user journeys',
        'Scope: Cart→Checkout→Payment→Tracking',
        'Run: On PR to staging branch',
        'Gate: Critical path failure = BLOCKED',
    ]),
    ('Layer 4\nPerformance', ORANGE, Inches(7.8), [
        'Framework: k6 + Lighthouse CI',
        'Target: LCP < 2.5s, API < 200ms p95',
        'Scope: Load test (100 concurrent users)',
        'Run: Pre-release to staging',
        'Gate: p95 > 500ms = BLOCKED',
    ]),
    ('Layer 5\nSecurity', ACCENT_RED, Inches(10.3), [
        'Tools: OWASP ZAP + Snyk + Dependabot',
        'Scope: XSS, SQLi, CSRF, dependency scan',
        'Run: Weekly + on production release',
        'Gate: Critical vuln = BLOCKED',
        'Snyk Team: $25/mo (production)',
    ]),
]
for title, color, left, items in test_layers:
    add_architecture_box(s, left, Inches(1.5), Inches(2.3), Inches(5.5), title, items, color)

add_textbox(s, Inches(0.3), Inches(0.9), Inches(12), Inches(0.5),
            'CI/CD Gate: Unit → Integration → E2E → Performance → Security → Staging → Production',
            font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════
# SLIDE 15: QUALITY GUARDRAIL VERIFICATION
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Quality Guardrail Verification — All 5 Guardrails ✅ PASS',
                   'Automated quality checks enforced by Supervisor Agent 00')

guardrails = [
    ('✅ Multi-Layer Testing', ACCENT_GREEN, [
        '5 testing layers defined (Unit → Integration → E2E → Perf → Security)',
        'CI/CD gates with blocker criteria for each layer',
        '80%+ unit test coverage threshold enforced',
        'Playwright E2E for critical user journeys',
    ]),
    ('✅ API Gateway Rate Limiting', ACCENT_GREEN, [
        'Token Bucket algorithm with per-endpoint rate limits',
        '12 endpoint groups with specific rate limits',
        'Next.js Middleware + Zod validation implementation',
        'POC: basic rate limiting via middleware',
    ]),
    ('✅ Infra Cost Dynamic Scaling', ACCENT_GREEN, [
        '5 scaling thresholds: 8→50→500→2,000→5,000+ merchants',
        'Cost scales progressively from ₹0 → ₹32,465 → ₹2,50,000+',
        'Auto-scaling triggers: CPU > 70%, Memory > 70%, Latency > 500ms',
        'Each threshold has mapped infra tier and services',
    ]),
    ('✅ POC Path Mapped to Full Arch', ACCENT_GREEN, [
        'POC = 48 of 103 requirements (47%), clearly documented',
        'Each built component linked to full-architecture container',
        'Each deferred component has deferral reason + target tier',
        'POC deployment topology built from same components',
    ]),
    ('✅ Architectural Diagrams Produced', ACCENT_GREEN, [
        'System Context (C4 L1), Container (C4 L2), Data Flow',
        'Deployment Topology, POC Scope, Google Stitch Pipeline',
        'WhatsApp Integration, Consolidated Delivery, Component Matrix',
        'All diagrams in Mermaid.js — viewable on GitHub/Notion',
    ]),
]
for i, (title, color, items) in enumerate(guardrails):
    y = Inches(1.2) + Inches(i * 1.2)
    add_shape_rect(s, Inches(0.3), y, Inches(12.7), Inches(1.0), WHITE)
    add_shape_rect(s, Inches(0.3), y, Inches(0.08), Inches(1.0), color)
    add_textbox(s, Inches(0.5), y + Inches(0.05), Inches(4), Inches(0.35),
                title, font_size=12, bold=True, color=color)
    add_multiline_textbox(s, Inches(4.5), y + Inches(0.05), Inches(8.5), Inches(0.9),
                          items, font_size=9, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 16: SCALING STRATEGY
# ════════════════════════════════════════
s = new_slide()
add_bg(s)
add_section_header(s, 'Scaling Strategy — From 8 Merchants to 5,000+',
                   'Progressive infrastructure scaling with clear cost and milestone thresholds')

scale_headers = ['Threshold', 'Merchants', 'DB Size', 'Infra Tier', 'Monthly Cost', 'Key Services']
scale_rows = [
    ['POC', '8', '~10MB', 'Supabase Free + Vercel Hobby', '₹0', 'All free tiers; no Redis/Meilisearch'],
    ['Beta', '50', '~100MB', 'Supabase Pro + Vercel Pro', '₹5,500', 'Add Redis ($9); PgFTS sufficient'],
    ['Growth', '500', '~1GB', 'Supabase Pro + Vercel Pro', '₹12,500', 'Add Meilisearch ($29); Sentry Team'],
    ['Scale', '2,000', '~4GB', 'Supabase Pro + Vercel Pro ×2', '₹32,465', 'Full monitoring; ShipRocket active'],
    ['Enterprise', '5,000+', '~10GB', 'Supabase Enterprise + Vercel Enterprise', '₹2,50,000+', 'Multi-region; DR; dedicated support'],
]
add_table(s, Inches(0.3), Inches(1.3), Inches(12.7), Inches(2.8), scale_headers, scale_rows,
          col_widths=[Inches(1.2), Inches(1.2), Inches(1.2), Inches(2.8), Inches(1.5), Inches(4.8)])

add_multiline_textbox(s, Inches(0.3), Inches(4.3), Inches(6), Inches(2.5), [
    'Auto-Scaling Triggers:',
    '',
    '• CPU > 70% → Scale up compute tier',
    '• Memory > 70% → Scale up database',
    '• API Latency > 500ms p95 → Add caching',
    '• Daily Active Users > 10K → Review arch',
    '• Error Rate > 1% → Investigate and rollback',
    '',
    'POC to Production Migration Path:',
    '1. Start: Supabase Free + Vercel Hobby (₹0)',
    '2. At 50 merchants: Upgrade to Pro tiers',
    '3. At 500: Add Meilisearch + enhanced monitoring',
    '4. At 2,000: Multi-region readiness + DR',
], font_size=10, color=DARK_TEXT)

add_multiline_textbox(s, Inches(6.8), Inches(4.3), Inches(6), Inches(2.5), [
    'Multi-City Expansion Architecture:',
    '',
    'Each new city = new market cluster:',
    '• Add city to geo_hierarchy table',
    '• Markets created under city (DB-driven)',
    '• Merchants register under market',
    '• Delivery zones created per city',
    '• Separate micro-hubs per cluster',
    '',
    'No code changes needed for cities:',
    'All configuration is database-driven.',
    'City selector in UI reads from DB.',
    '',
    'Target Cities (Year 2):',
    'Mysore, Mangalore, Hubli, Davangere',
    '→ 5,000+ merchants across Karnataka',
], font_size=10, color=DARK_TEXT)

# ════════════════════════════════════════
# SLIDE 17: CLOSING
# ════════════════════════════════════════
s = new_slide()
add_bg(s, BG_DARK)
add_textbox(s, Inches(0.8), Inches(1.5), Inches(11), Inches(1),
            'PeteMart', font_size=54, bold=True, color=ACCENT_GOLD)
add_textbox(s, Inches(0.8), Inches(2.7), Inches(11), Inches(0.8),
            'Enterprise Architecture Blueprint — Ready for Review',
            font_size=18, color=WHITE)
add_shape_rect(s, Inches(0.8), Inches(3.7), Inches(3), Inches(0.04), ACCENT_GOLD)
add_multiline_textbox(s, Inches(0.8), Inches(4.0), Inches(11), Inches(2.5), [
    'Dual Architecture: Full Product (103 reqs) + POC (48 reqs, ₹0/month)',
    'All 5 Quality Guardrails: ✅ PASS | All Diagrams: ✅ Produced',
    'POC: 9 merchants · 15 tables · 3 modes · Zero-cost · 6-8 weeks',
    'Production: ₹32,465/mo cloud | Scale to 5,000+ merchants across Karnataka',
    '',
    'Next Phase: Agent 04 (Prototype) builds the POC workspace',
    'Awaiting GATE-TECH-STACK-01 & GATE-COSTING-01 approval',
], font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

add_textbox(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
            'Thank You | Questions & Architecture Review', font_size=16, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

# ── Save ──
output_path = os.path.join(os.getcwd(), 'agents', '03_architect_agent', '03_architect_agent_FULL_PRESENTATION.pptx')
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
